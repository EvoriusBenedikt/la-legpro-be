from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, BackgroundTasks
import sqlite3, os, json, time, re
from typing import List, Optional, Any
from pydantic import BaseModel
import auth
import chromadb
from services.llm_client import call_glm

router = APIRouter(prefix="/api", tags=["compliance"])

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

class AnalyzeRequest(BaseModel):
    doc_id: int

class PasalAnalyzeRequest(BaseModel):
    doc_id: int
    pasals: List[str]

@router.post("/upload")
async def upload_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...), 
    doc_type: str = Form("regulations"),
    klasifikasi: str = Form("Umum"),
    current_user: dict = Depends(auth.require_role("manajer"))
):
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
        
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
    # check file size without reading into memory completely if possible
    file.file.seek(0, 2)
    file_size = file.file.tell()
    file.file.seek(0)
    
    if file_size > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File too large (max 50MB)")
        
    import sqlite3
    import uuid
    from datetime import datetime
    
    # Setup directories
    pdfs_dir = os.path.join(BASE_DIR, "data", "pdfs")
    os.makedirs(pdfs_dir, exist_ok=True)
    
    file_path = os.path.join(pdfs_dir, file.filename)
    
    # Save physical file temporarily
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        # ── Proceed with Saving Metadata ─────────────────────────────────────
        db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
        conn = sqlite3.connect(db_path, timeout=30.0)
        c = conn.cursor()
        c.execute('''CREATE TABLE IF NOT EXISTS regulations
                     (id INTEGER PRIMARY KEY AUTOINCREMENT, domain TEXT, judul TEXT, 
                      nomor TEXT, jenis TEXT, sektor TEXT, status TEXT, 
                      detail_url TEXT, download_url TEXT, local_path TEXT, klasifikasi TEXT)''')
                      
        # Generate ID and Metadata
        doc_id = str(uuid.uuid4())[:8]
        judul = file.filename.replace('.pdf', '')
        nomor = f"CUSTOM-{doc_id}"
        
        if doc_type == "internal":
            jenis = "Dokumen Internal"
            sektor = "Dokumen Internal"
        else:
            jenis = "Regulasi Custom"
            sektor = "Upload Manual"
        
        status = "Memproses" # Initialize with processing status
        
        c.execute('''INSERT OR REPLACE INTO regulations 
                     (judul, download_url, nomor, jenis, sektor, status, detail_url, local_path, klasifikasi, domain) 
                     VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)''',
                  (judul, "", nomor, jenis, sektor, status, "", file_path, klasifikasi, "Custom"))
        # Get the actual auto-incremented ID
        doc_id = str(c.lastrowid)
        conn.commit()
        conn.close()
        
        # Enqueue background processing task
        background_tasks.add_task(process_document_background, file_path, doc_id, file.filename, nomor, jenis, sektor, status, klasifikasi)
        
        return {"status": "success", "message": "Dokumen berhasil diunggah dan sedang diproses di latar belakang."}
        
    except Exception as e:
        print(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail="Terjadi kesalahan saat memproses dokumen.")

@router.post("/analyze")
async def analyze_document(req: AnalyzeRequest):
    """Analyzes a document from ChromaDB chunks + LLM for overview and status."""
    collection = get_chroma_collection()
    
    # Step 1: Fetch all chunks for this specific document
    # Prefer reg_id (unique) but fall back to nomor if not available
    try:
        if req.reg_id:
            results = collection.get(
                where={"reg_id": req.reg_id},
                include=["documents"]
            )
        else:
            results = collection.get(
                where={"nomor": req.nomor},
                include=["documents"]
            )
        chunks = results.get("documents", [])
    except Exception as e:
        print(f"ChromaDB fetch error: {e}")
        chunks = []
    
    if not chunks:
        return {
            "total_pasal": 0,
            "overview": "Dokumen ini belum memiliki data di knowledge base.",
            "status": {"dicabut": [], "diubah_dengan": []},
            "outline": []
        }
    
    # Step 2: Build outline + find max Pasal number (more accurate than chunk count)
    import re
    outline = []
    seen = set()
    max_pasal_num = 0
    
    for chunk in chunks:
        lines = chunk.strip().split('\n')
        first_line = lines[0].strip() if lines else ""
        
        # Detect BAB headers
        if re.match(r'^BAB\s+[IVXLC\d]+', first_line, re.IGNORECASE):
            entry = first_line[:80]
            if entry not in seen:
                outline.append({"type": "bab", "text": entry})
                seen.add(entry)
        
        # Detect Pasal headers — track the max number found
        pasal_match = re.search(r'Pasal\s+(\d+)', chunk, re.IGNORECASE)
        if pasal_match:
            num = int(pasal_match.group(1))
            if num > max_pasal_num:
                max_pasal_num = num
            
            entry = first_line[:60]
            if re.match(r'^Pasal\s+\d+', first_line, re.IGNORECASE) and entry not in seen:
                outline.append({"type": "pasal", "text": entry})
                seen.add(entry)

    # Step 3: Take sample text for LLM (first 6 chunks ≈ 3500 chars)
    sample_text = "\n---\n".join(chunks[:6])[:3500]
    
    system_prompt = (
        "Anda adalah analis hukum OJK Indonesia. Berikan analisis singkat dan akurat dalam format JSON yang ketat. "
        "Jangan menambahkan teks di luar JSON."
    )
    
    user_prompt = (
        f"Analisis dokumen hukum berikut:\n"
        f"Judul: {req.judul}\nNomor: {req.nomor}\n\n"
        f"TEKS DOKUMEN:\n{sample_text}\n\n"
        f"Kembalikan HANYA JSON berikut (tanpa markdown, tanpa kode blok):\n"
        '{"overview": "ringkasan 2-3 kalimat", '
        '"dicabut": ["nama peraturan yang dicabut jika ada, atau array kosong"], '
        '"diubah_dengan": ["nama peraturan pengubah jika ada, atau array kosong"]}'
    )
    
    messages_llm = [
        {"role": "system", "content": system_prompt},
        {"role": "user",   "content": user_prompt}
    ]
    
    overview = "Tidak dapat mengambil ringkasan dari dokumen ini."
    dicabut = []
    diubah_dengan = []
    
    try:
        raw_content = call_glm(messages_llm, temperature=0.1, timeout=60)
        raw_content = re.sub(r'```json|```', '', raw_content).strip()
        parsed = json.loads(raw_content)
        overview = parsed.get("overview", overview)
        dicabut = parsed.get("dicabut", [])
        diubah_dengan = parsed.get("diubah_dengan", [])
    except HTTPException:
        raise
    except Exception as e:
        print(f"LLM analysis error: {e}")
    
    return {
        "total_pasal": max_pasal_num,
        "overview": overview,
        "status": {"dicabut": dicabut, "diubah_dengan": diubah_dengan},
        "outline": outline
    }

@router.post("/analyze-pasals")
async def analyze_pasals(req: AnalyzeRequest):
    """
    Deep analysis: uses LLM to analyze each Pasal individually,
    producing status, perbandingan (comparison), and hubungan (hierarchy).
    """
    collection = get_chroma_collection()

    try:
        if req.reg_id:
            results = collection.get(where={"reg_id": req.reg_id}, include=["documents"])
        else:
            results = collection.get(where={"nomor": req.nomor}, include=["documents"])
        chunks = results.get("documents", [])
    except Exception as e:
        print(f"ChromaDB error: {e}")
        chunks = []

    if not chunks:
        return {"pasals": [], "error": "Tidak ada data untuk dokumen ini."}

    # ── Build pasal map: {pasal_num (int): text (str)} ──
    pasal_texts: dict[int, str] = {}
    for chunk in chunks:
        m = re.search(r'Pasal\s+(\d+)', chunk, re.IGNORECASE)
        if m:
            num = int(m.group(1))
            if num not in pasal_texts:
                pasal_texts[num] = chunk[:500]

    sorted_pasals = sorted(pasal_texts.items())[:15]   # limit to 15 most relevant pasals

    if not sorted_pasals:
        return {"pasals": [], "error": "Tidak dapat mengekstrak Pasal dari dokumen ini."}

    # ── Build LLM prompt ──
    doc_text = ""
    for num, text in sorted_pasals:
        doc_text += f"\nPasal {num}:\n{text[:380]}\n---\n"

    system_prompt = (
        "Anda adalah analis regulasi di Indonesia yang sangat teliti dan akurat. "
        "Selalu jawab dalam Bahasa Indonesia. Jangan menambahkan teks di luar JSON."
    )
    user_prompt = (
        f"Analisis regulasi berikut secara mendalam:\n"
        f"Judul: {req.judul}\nNomor: {req.nomor}\n\n"
        f"TEKS REGULASI:\n{doc_text}\n\n"
        f"Untuk SETIAP Pasal di atas, berikan analisis dalam format JSON strict berikut "
        f"(kembalikan HANYA array JSON, tanpa teks tambahan):\n"
        '[{"pasal":"Pasal 1","status":"Tidak Berubah",'
        '"isi":"ringkasan isi pasal dalam 1-2 kalimat",'
        '"perbandingan":"perbandingan konten pasal ini dengan regulasi sebelumnya yang dicabut atau diubah",'
        '"hubungan":"hubungan dengan Undang-Undang atau Peraturan Pemerintah yang lebih tinggi"},'
        '{"pasal":"Pasal 2",...}]'
        '\n\nNilai status yang valid: "Tidak Berubah", "Diubah", "Baru", "Dicabut".'
    )

    try:
        raw = call_glm([
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_prompt},
        ], temperature=0.1, timeout=180)
        raw = re.sub(r'```json|```', '', raw).strip()

        # The LLM might return an object {"pasals":[...]} OR a bare array [...]
        parsed = json.loads(raw)
        if isinstance(parsed, dict):
            items = parsed.get("pasals", parsed.get("data", []))
        elif isinstance(parsed, list):
            items = parsed
        else:
            items = []

        # Attach the actual extracted pasal text for the frontend to display
        pasal_map = {num: text for num, text in sorted_pasals}
        for item in items:
            m = re.search(r'\d+', item.get("pasal", ""))
            if m:
                item["isi_teks"] = pasal_map.get(int(m.group()), "")[:350]

        return {"pasals": items}

    except json.JSONDecodeError as e:
        print(f"JSON parse error in analyze-pasals: {e}\nRaw: {raw[:300]}")
        return {"pasals": [], "error": "LLM memberikan respons yang tidak valid. Coba lagi."}
    except Exception as e:
        print(f"Deep analysis error: {e}")
        return {"pasals": [], "error": str(e)}


def process_single_pasal(c, doc_type, pihak_1, pihak_2, sektor, pokok, collection, COSINE_THRESHOLD):
    pasal_label = c.get("pasal", "Pasal")
    desc = c.get("deskripsi", "")
    if not desc or len(desc) < 10:
        return None

    # Map doc_type to doc_category
    doc_cat = "NDA" if "NDA" in str(doc_type).upper() or "NON-DISCLOSURE" in str(doc_type).upper() else "PKS"

    # Enrich the search query with document context for better RAG retrieval
    search_query = f"[{doc_type}] [{sektor}] {desc[:400]}"
    sr = collection.query(
        query_texts=[search_query], 
        n_results=3,
        where={"doc_category": {"$in": [doc_cat, "UMUM"]}}
    )

    supporting_regulations = []
    relevant_refs = []

    if sr and sr.get('documents') and len(sr['documents'][0]) > 0:
        for i in range(len(sr['documents'][0])):
            dist = sr['distances'][0][i]
            if dist < COSINE_THRESHOLD:
                doc_text = sr['documents'][0][i]
                meta = sr['metadatas'][0][i]
                reg_name = f"{meta.get('jenis', 'Aturan')} Nomor {meta.get('nomor', 'N/A')}"

                # Relevance confirmation micro-call
                if is_regulation_relevant(desc, doc_text, reg_name):
                    supporting_regulations.append({
                        "jenis": meta.get('jenis', 'Aturan'),
                        "nomor": meta.get('nomor', 'N/A'),
                        "sektor": meta.get('sektor', 'Umum'),
                        "teks": doc_text
                    })
                    relevant_refs.append(f"--- {reg_name} ---\n{doc_text[:700]}")
                else:
                    print(f"  [Relevance] {reg_name} rejected for {pasal_label}")

    if not relevant_refs:
        return {
            "pasal": pasal_label,
            "isi_pasal": desc[:500],
            "status": "TIDAK DIATUR",
            "penjelasan": "Tidak ditemukan regulasi yang secara langsung relevan dengan klausul ini dalam database saat ini.",
            "rekomendasi": "Pastikan klausul ini wajar secara keperdataan dan tidak bertentangan dengan KUH Perdata.",
            "ai_analysis": "",
            "supporting_regulations": supporting_regulations
        }

    combined_regs = "\n\n".join(relevant_refs)
    prompt_verify = (
        "Anda adalah auditor kepatuhan hukum senior Indonesia yang sangat teliti. "
        "Analisis klausul kontrak berikut secara mendalam menggunakan regulasi yang tersedia.\n\n"

        f"KONTEKS DOKUMEN:\n"
        f"- Jenis: {doc_type}\n"
        f"- Pihak Pertama: {pihak_1}\n"
        f"- Pihak Kedua: {pihak_2}\n"
        f"- Sektor: {sektor}\n"
        f"- Pokok Perjanjian: {pokok}\n\n"

        f"KLAUSUL YANG DIANALISIS ({pasal_label}):\n{desc}\n\n"

        f"REGULASI YANG BERLAKU:\n{combined_regs}\n\n"

        "INSTRUKSI ANALISIS (ikuti urutan ini):\n"
        "1. Identifikasi: Apa topik hukum utama klausul ini? (pembayaran, kerahasiaan, PHK, dll)\n"
        "2. Kesesuaian: Apakah klausul ini SECARA EKSPLISIT diatur, dilarang, atau diwajibkan oleh regulasi di atas?\n"
        "   Kutip pasal/ayat spesifik dari regulasi jika ada.\n"
        "3. Risiko: Identifikasi risiko hukum konkret bagi pihak yang lebih lemah.\n"
        "4. Verdict: Berikan satu dari tiga status:\n"
        "   - SESUAI: klausul patuh dan tidak berisiko\n"
        "   - BERESIKO: klausul perlu perbaikan atau klarifikasi\n"
        "   - FATAL: klausul melanggar ketentuan wajib regulasi\n\n"
        "Kembalikan HANYA JSON berikut (tanpa markdown, tanpa teks lain):\n"
        '{\n'
        '  "status": "SESUAI" | "BERESIKO" | "FATAL",\n'
        '  "pasal_regulasi_terkait": "misal: Pasal 15 ayat (3) PP 45/2015",\n'
        '  "penjelasan": "penjelasan spesifik 1-2 kalimat merujuk regulasi",\n'
        '  "rekomendasi": "saran perbaikan konkret (kosongkan jika SESUAI)",\n'
        '  "ai_analysis": "dampak bisnis jangka panjang (isi hanya jika BERESIKO atau FATAL)"\n'
        '}'
    )

    try:
        raw_v = call_glm(
            [{"role": "user", "content": prompt_verify}],
            temperature=0.1,
            timeout=90
        )
        raw_v = re.sub(r'```json|```', '', raw_v).strip()
        ans = json.loads(raw_v)

        status_raw = str(ans.get("status", "BERESIKO")).upper()
        if "SESUAI" in status_raw:       status_final = "SESUAI"
        elif "FATAL" in status_raw:       status_final = "FATAL"
        elif "TIDAK DIATUR" in status_raw: status_final = "TIDAK DIATUR"
        else:                              status_final = "BERESIKO"

        return {
            "pasal": pasal_label,
            "isi_pasal": desc[:500],
            "status": status_final,
            "pasal_regulasi_terkait": ans.get("pasal_regulasi_terkait", ""),
            "penjelasan": ans.get("penjelasan", "Gagal diverifikasi"),
            "rekomendasi": ans.get("rekomendasi", ""),
            "ai_analysis": ans.get("ai_analysis", ""),
            "supporting_regulations": supporting_regulations
        }

    except HTTPException as e:
        print(f"[Compliance] Verification HTTP error: {e.detail}")
        return {
            "pasal": pasal_label, "isi_pasal": desc[:500], "status": "ERROR",
            "pasal_regulasi_terkait": "",
            "penjelasan": "Layanan AI sedang mengalami gangguan koneksi sementara. Silakan coba lagi dalam beberapa saat.",
            "rekomendasi": "Coba ulang analisis.",
            "ai_analysis": "", "supporting_regulations": supporting_regulations
        }
    except Exception as e:
        print(f"[Compliance] Parsing error for {pasal_label}: {e}")
        return {
            "pasal": pasal_label, "isi_pasal": desc[:500], "status": "ERROR",
            "pasal_regulasi_terkait": "",
            "penjelasan": "Format respons LLM tidak valid.",
            "rekomendasi": "-",
            "ai_analysis": "", "supporting_regulations": supporting_regulations
        }

def extract_text_multi_format(file_path: str, filename: str, use_ocr: bool = False) -> str:
    """
    Extracts text from various file formats: PDF, DOCX, XLSX, PPTX, JPG, PNG, TXT.
    Routes to the appropriate extractor based on the file extension.
    """
    ext = filename.lower().split('.')[-1]
    
    if ext == 'pdf':
        if use_ocr:
            print(f"  [OCR] Forcing traditional OCR extraction for {filename}")
            import sys
            if BASE_DIR not in sys.path:
                sys.path.append(BASE_DIR)
            from parser.pdf_parser import LegalDocumentParser
            parser = LegalDocumentParser()
            return parser.parse_pdf(file_path, force_ocr=True)
        else:
            return extract_text_hybrid(file_path, force_vlm=False)
        
    elif ext == 'txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
            
    elif ext in ['jpg', 'jpeg', 'png']:
        if use_ocr:
            print(f"  [OCR] Traditional image extraction for {filename}")
            import sys
            if BASE_DIR not in sys.path:
                sys.path.append(BASE_DIR)
            from parser.pdf_parser import LegalDocumentParser
            parser = LegalDocumentParser()
            if parser.ocr:
                import numpy as np
                from PIL import Image
                img = Image.open(file_path).convert("RGB")
                img_array = np.array(img)
                result = parser.ocr.ocr(img_array, cls=False)
                page_text = []
                if result and result[0]:
                    for line in result[0]:
                        page_text.append(line[1][0])
                return " ".join(page_text)
            else:
                return "[OCR GAGAL INISIALISASI]"
        else:
            print(f"  [VLM] Image extraction for {filename}")
            with open(file_path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode('utf-8')
            vlm_text = call_glm_vision(img_b64, VLM_PAGE_PROMPT, timeout=60)
            return vlm_text.strip()
        
    elif ext == 'docx':
        try:
            import docx
            doc = docx.Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:
                        full_text.append(" | ".join(row_data))
            return "\n".join(full_text)
        except Exception as e:
            print(f"Error parsing DOCX: {e}")
            return ""
            
    elif ext == 'pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        full_text.append(shape.text.strip())
            return "\n".join(full_text)
        except Exception as e:
            print(f"Error parsing PPTX: {e}")
            return ""
            
    elif ext == 'xlsx':
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            ws = wb.active
            full_text = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if row_data:
                    full_text.append(" | ".join(row_data))
            return "\n".join(full_text)
        except Exception as e:
            print(f"Error parsing XLSX: {e}")
            return ""
            
    else:
        return ""


@router.post("/check-compliance")
async def check_compliance(file: UploadFile = File(...), use_ocr: str = Form("false")):
    """
    Accepts a document upload (PDF, DOCX, XLSX, PPTX, JPG, PNG, TXT), extracts text, 
    and uses a multi-step LLM pipeline to cross-check clauses against the OJK repository.
    """
    allowed_exts = ['.pdf', '.docx', '.xlsx', '.pptx', '.jpg', '.jpeg', '.png', '.txt']
    ext = os.path.splitext(file.filename.lower())[1]
    
    if ext not in allowed_exts:
        raise HTTPException(status_code=400, detail="Format file tidak didukung. Harap unggah PDF, DOCX, XLSX, PPTX, JPG, PNG, atau TXT.")
        
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
    file.file.seek(0, 2)
    file_size = file.file.tell()
    file.file.seek(0)
    
    if file_size > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File too large (max 50MB)")
        
    upload_id = str(uuid.uuid4())[:8]
    temp_path = os.path.join(PDFS_DIR, f"temp_check_{upload_id}_{file.filename}")
    
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        # ── Pass 0: Multi-Format Extraction ──────────────────────────────────
        print(f"[Compliance] Starting extraction for: {file.filename}")
        is_ocr = use_ocr.lower() == "true"
        full_text = extract_text_multi_format(temp_path, file.filename, use_ocr=is_ocr)
        text_snippet = full_text[:35000]

        if len(text_snippet.strip()) < 100:
            raise HTTPException(
                status_code=400,
                detail="Dokumen tidak dapat dibaca. Coba unggah ulang atau pastikan PDF tidak terenkripsi."
            )

        # ── Pass 0A: Document Metadata Understanding ────────────────────────
        print("[Compliance] Pass 0A: Extracting document metadata...")
        doc_ctx = understand_document(full_text)
        doc_type = doc_ctx.get("jenis_dokumen", "Kontrak")
        pihak_1  = doc_ctx.get("pihak_pertama", "Pihak Pertama")
        pihak_2  = doc_ctx.get("pihak_kedua", "Pihak Kedua")
        sektor   = doc_ctx.get("sektor_bisnis", "umum")
        pokok    = doc_ctx.get("pokok_perjanjian", "")

        # ── Pass 0B: Dedicated signing date extraction ──────────────────────
        print("[Compliance] Pass 0B: Extracting signing date...")
        tanggal_pembuatan = extract_signing_date(full_text)

        # ── Pass 0C: Dedicated duration extraction ──────────────────────────
        print("[Compliance] Pass 0C: Extracting contract duration...")
        durasi_bulan = extract_duration_months(full_text)

        # ── Compute expiry date ─────────────────────────────────────────────
        tanggal_berakhir = None
        if tanggal_pembuatan and durasi_bulan:
            try:
                start_date = datetime.strptime(tanggal_pembuatan, "%Y-%m-%d")
                end_date = start_date + relativedelta(months=int(durasi_bulan))
                tanggal_berakhir = end_date.strftime("%Y-%m-%d")
                print(f"[Compliance] Expiry computed: {tanggal_pembuatan} + {durasi_bulan}mo = {tanggal_berakhir}")
            except Exception as e:
                print(f"[Compliance] Failed to compute expiry: {e}")
                
        # ── Pass 0D: Explicit end date fallback ─────────────────────────────
        if not tanggal_berakhir:
            print("[Compliance] Computation failed/missing, attempting explicit end date extraction...")
            tanggal_berakhir = extract_explicit_end_date(full_text)

        # ── Pass 1: Clause Extraction ───────────────────────────────────────
        print("[Compliance] Pass 1: Extracting clauses...")
        pasal_items = extract_pasal_items(text_snippet, max_items=20, max_chars_per_pasal=1200)

        # Fallback to LLM extractor for non-standard contracts
        if len(pasal_items) < 2:
            print("[Compliance] Regex found <2 clauses, falling back to LLM extractor...")
            pasal_items = llm_extract_pasal_items(full_text, max_items=20)

        # Final fallback: treat whole document as one unit
        if not pasal_items:
            pasal_items = [{"pasal": "Dokumen Keseluruhan", "deskripsi": text_snippet[:1200]}]

        print(f"[Compliance] Found {len(pasal_items)} clauses. Running analysis...")

        # ── Pass 2: Per-clause RAG + LLM Compliance Check ──────────────────
        collection = get_chroma_collection()
        results = []
        COSINE_THRESHOLD = 0.45  # tightened from 0.60 to reduce wrong-domain citations

        def worker(c):
            return process_single_pasal(c, doc_type, pihak_1, pihak_2, sektor, pokok, collection, COSINE_THRESHOLD)

        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            for res in executor.map(worker, pasal_items[:20]):
                if res:
                    results.append(res)

        # ── Pass 4: Document-Level Summary ─────────────────────────────────
        fatal_count    = sum(1 for r in results if r["status"] == "FATAL")
        beresiko_count = sum(1 for r in results if r["status"] == "BERESIKO")
        sesuai_count   = sum(1 for r in results if r["status"] == "SESUAI")
        tidak_diatur_count = sum(1 for r in results if r["status"] == "TIDAK DIATUR")

        # Calculate score only against regulated clauses
        regulated_count = len(results) - tidak_diatur_count
        skor = round((sesuai_count / max(regulated_count, 1)) * 100) if regulated_count > 0 else 100

        summary = {
            "jenis_dokumen": doc_type,
            "pihak_pertama": pihak_1,
            "pihak_kedua": pihak_2,
            "sektor_bisnis": sektor,
            "pokok_perjanjian": pokok,
            "tanggal_berakhir": tanggal_berakhir,
            "total_klausul": len(results),
            "sesuai": sesuai_count,
            "beresiko": beresiko_count,
            "fatal": fatal_count,
            "skor_kepatuhan": skor
        }

        print(f"[Compliance] Done. Score: {summary['skor_kepatuhan']}% | "
              f"SESUAI={sesuai_count} BERESIKO={beresiko_count} FATAL={fatal_count}")

        return {"report": results, "summary": summary}

    except HTTPException:
        raise
    except Exception as e:
        print(f"Compliance checker unexpected error: {e}")
        raise HTTPException(status_code=500, detail="Terjadi kesalahan internal saat memproses dokumen.")
    finally:
        # Automagically clean up uploaded PKS memory footprint
        if os.path.exists(temp_path):
            try: os.remove(temp_path)
            except: pass

class KGExclusionCreate(BaseModel):
    entity_name: str
    import uvicorn
    # Natively running on port 8000
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=False)
