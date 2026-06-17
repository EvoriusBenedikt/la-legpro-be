import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from dotenv import load_dotenv
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
load_dotenv(os.path.join(BASE_DIR, ".env"))

import re
import json
import time
import requests
import chromadb
from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Depends, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel, field_validator
import auth
from typing import List, Optional, Any
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
import shutil
import hashlib
import uuid
import base64
import logging
import concurrent.futures
from datetime import datetime
from dateutil.relativedelta import relativedelta

# ── File logging setup ──────────────────────────────────────────────────────
_LOG_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "server.log")
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(message)s",
    handlers=[
        logging.FileHandler(_LOG_PATH, encoding="utf-8"),
        logging.StreamHandler(sys.stdout),
    ],
)
# Redirect all print() calls to the logger so nothing is missed
import builtins as _builtins
_real_print = _builtins.print
def _log_print(*args, **kwargs):
    msg = " ".join(str(a) for a in args)
    logging.info(msg)
_builtins.print = _log_print


# ── Load environment ────────────────────────────────────────────────────────

CHROMA_DB_DIR = os.path.join(BASE_DIR, "data", "chroma_db")

# ── GLM API Config ──────────────────────────────────────────────────────────
API_BASE_URL = os.getenv("MODEL_BASE_URL", "https://console.labahasa.ai/v1").rstrip("/")
API_KEY      = os.getenv("MODEL_API_KEY", "")
# Reverted to use Maverick for both text and vision because GLM is unstable
GLM_MODEL    = os.getenv("LLAMA_MODEL", "llama-4-maverick-instruct")
VLM_MODEL    = os.getenv("LLAMA_MODEL", "llama-4-maverick-instruct")
GLM_MAX_ATTEMPTS = int(os.getenv("GLM_MAX_ATTEMPTS", "5"))
GLM_RETRY_BACKOFF_BASE = float(os.getenv("GLM_RETRY_BACKOFF_BASE", "0.8"))

if not API_BASE_URL or not API_KEY:
    print("[WARNING] MODEL_BASE_URL or MODEL_API_KEY is not set in .env — LLM calls will fail.")

# Setup Paths for Parser
sys.path.append(os.path.join(BASE_DIR, "parser"))
from pdf_parser import LegalDocumentParser, LegalChunker

# Setup FastAPI App
app = FastAPI(title="Legal Analyzer API")

# ── Init DB for metadata ────────────────────────────────────────────────────
def init_main_db():
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    os.makedirs(os.path.dirname(db_path), exist_ok=True)
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS regulations
                 (id TEXT PRIMARY KEY, judul TEXT, pdf_url TEXT, 
                  nomor TEXT, jenis TEXT, sektor TEXT, 
                  status TEXT, detail_url TEXT, local_path TEXT)''')
    try:
        c.execute("ALTER TABLE regulations ADD COLUMN klasifikasi TEXT DEFAULT 'Umum'")
    except sqlite3.OperationalError:
        pass
        
    c.execute('''CREATE TABLE IF NOT EXISTS access_grants (
        id TEXT PRIMARY KEY,
        doc_id TEXT,
        granted_by TEXT,
        granted_to TEXT,
        reason TEXT,
        expires_at TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )''')
    c.execute('''CREATE TABLE IF NOT EXISTS audit_logs (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        user_id TEXT,
        action_type TEXT,
        resource_id TEXT,
        details TEXT
    )''')
    c.execute('''CREATE TABLE IF NOT EXISTS kg_nodes (
        id TEXT PRIMARY KEY,
        label TEXT NOT NULL,
        type TEXT NOT NULL,
        doc_id TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )''')
    c.execute('''CREATE TABLE IF NOT EXISTS kg_edges (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        source_id TEXT NOT NULL,
        target_id TEXT NOT NULL,
        relation TEXT NOT NULL,
        doc_id TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )''')
    
    c.execute('''CREATE TABLE IF NOT EXISTS kg_exclusions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        entity_name TEXT UNIQUE NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )''')
    
    # FR-29: Document Taxonomy
    c.execute('''CREATE TABLE IF NOT EXISTS document_taxonomy (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT UNIQUE NOT NULL,
        parent_id INTEGER,
        is_active BOOLEAN DEFAULT 1,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )''')
    
    # Seed default taxonomy if empty
    c.execute("SELECT COUNT(*) FROM document_taxonomy")
    if c.fetchone()[0] == 0:
        default_types = [
            "Peraturan Pemerintah",
            "Undang-Undang",
            "Peraturan OJK",
            "Surat Edaran OJK",
            "Dokumen Internal",
            "Peraturan Menteri",
            "Regulasi Custom"
        ]
        for dt in default_types:
            c.execute("INSERT INTO document_taxonomy (name) VALUES (?)", (dt,))

    
    conn.commit()
    conn.close()

def log_audit(user_id: str, action_type: str, resource_id: str = "", details: str = ""):
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    try:
        conn = sqlite3.connect(db_path, timeout=10.0)
        c = conn.cursor()
        c.execute('''INSERT INTO audit_logs (user_id, action_type, resource_id, details) 
                     VALUES (?, ?, ?, ?)''', (user_id, action_type, resource_id, details))
        conn.commit()
        conn.close()
    except Exception as e:
        print(f"Failed to write audit log: {e}")

init_main_db()

def extract_and_store_graph(doc_id: str, full_text: str, nomor: str, judul: str, jenis: str):
    """Use LLM to extract entities & relationships from a document, then upsert into KG tables."""
    import sqlite3, json as _json
    import os
    
    # FR-30: Fetch Exclusions
    exclusions = []
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    try:
        conn = sqlite3.connect(db_path, timeout=30.0)
        c = conn.cursor()
        c.execute("SELECT entity_name FROM kg_exclusions")
        exclusions = [row[0] for row in c.fetchall()]
        conn.close()
    except Exception as e:
        print(f"[KG] Warning: Failed to fetch exclusions: {e}")
        
    exclusion_text = ""
    if exclusions:
        exclusion_list_str = ", ".join(exclusions)
        exclusion_text = f"\n6. DILARANG KERAS mengekstrak entitas berikut ini (Abaikan mereka sepenuhnya): {exclusion_list_str}."

    snippet = full_text[:8000]  # Increased from 3500 to capture definitions (Pasal 1) and core body
    prompt = [
        {"role": "system", "content": (
            "Kamu adalah ekstraktor entitas hukum level ahli. Baca teks peraturan di bawah ini dan ekstrak "
            "entitas serta hubungan hukumnya ke dalam bentuk JSON murni. "
            "Format JSON yang diharapkan:\n"
            "{\"entitas\": [{\"id\": \"string unik\", \"label\": \"nama tampilan\", \"type\": \"entitas|topik\"}], "
            "\"relasi\": [{\"source\": \"id sumber\", \"target\": \"id target\", \"rel\": \"MENCABUT|MENGUBAH|MERUJUK|MENGATUR|DITERBITKAN_OLEH\"}]}\n"
            "Aturan Ekstraksi Kritis:\n"
            "1. Untuk type='entitas': Selalu ekstrak lembaga penerbit (e.g. OJK, Kemnaker, BI, Kemenkeu, Presiden).\n"
            "2. Untuk type='topik': Sangat penting untuk mendeteksi topik terkait skenario NDA dan PKS! Jika teks mengandung unsur 'Kerahasiaan', 'Data Pribadi', 'Keamanan Informasi', 'Rahasia Dagang', ekstrak sebagai topik NDA. Jika mengandung 'Perjanjian', 'Kontrak', 'Kemitraan', 'Vendor', 'Pengadaan', ekstrak sebagai topik PKS.\n"
            "3. Selalu buat relasi DITERBITKAN_OLEH dari regulasi ini ke lembaga penerbitnya.\n"
            "4. Gunakan nomor regulasi resmi (misal POJK-12-2023) sebagai ID untuk target relasi MERUJUK/MENGUBAH.\n"
            "5. Jangan batasi jumlah ekstraksi. Ekstrak SEMUA entitas, topik relevan, dan regulasi terkait yang ada dalam teks untuk membangun Knowledge Graph yang padat dan komprehensif."
            f"{exclusion_text}"
        )},
        {"role": "user", "content": f"Regulasi: {jenis} Nomor {nomor}\nJudul: {judul}\n\nTeks:\n{snippet}"}
    ]
    try:
        raw = call_glm(prompt, temperature=0.0, timeout=45)
        # Strip markdown code fences if present
        raw = raw.strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        data = _json.loads(raw)
        
        # FR-30 Post-processing: remove excluded entities
        if exclusions:
            exc_lower = {e.lower() for e in exclusions}
            # Filter nodes
            valid_entities = []
            excluded_ids = set()
            for ent in data.get("entitas", []):
                if ent.get("label", "").lower() in exc_lower:
                    excluded_ids.add(ent.get("id"))
                else:
                    valid_entities.append(ent)
            data["entitas"] = valid_entities
            
            # Filter edges referencing excluded nodes
            valid_relations = []
            for rel in data.get("relasi", []):
                if rel.get("source") not in excluded_ids and rel.get("target") not in excluded_ids:
                    valid_relations.append(rel)
            data["relasi"] = valid_relations

    except Exception as e:
        print(f"[KG] LLM extraction failed for {nomor}: {e}")
        return

    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    try:
        conn = sqlite3.connect(db_path, timeout=30.0)
        c = conn.cursor()

        # Upsert the regulation itself as a node
        reg_node_id = f"{jenis}-{nomor}".replace(" ", "-")[:80]
        c.execute("INSERT OR REPLACE INTO kg_nodes (id, label, type, doc_id) VALUES (?, ?, ?, ?)",
                  (reg_node_id, f"{jenis} {nomor}", "regulasi", doc_id))

        # Upsert extracted entities/topics
        for ent in data.get("entitas", []):
            ent_id = str(ent.get("id", "")).strip()[:80]
            ent_label = str(ent.get("label", ent_id)).strip()[:120]
            ent_type = str(ent.get("type", "topik")).strip()
            if ent_id:
                c.execute("INSERT OR IGNORE INTO kg_nodes (id, label, type, doc_id) VALUES (?, ?, ?, ?)",
                          (ent_id, ent_label, ent_type, None))

        # Insert edges
        for rel in data.get("relasi", []):
            src = str(rel.get("source", "")).strip()[:80]
            tgt = str(rel.get("target", "")).strip()[:80]
            relation = str(rel.get("rel", "MERUJUK")).strip()[:40]
            if src and tgt and src != tgt:
                # Replace regulation self-reference with the canonical reg_node_id
                if src == nomor or src == f"{jenis} {nomor}":
                    src = reg_node_id
                if tgt == nomor or tgt == f"{jenis} {nomor}":
                    tgt = reg_node_id
                c.execute("INSERT INTO kg_edges (source_id, target_id, relation, doc_id) VALUES (?, ?, ?, ?)",
                          (src, tgt, relation, doc_id))

        conn.commit()
        conn.close()
        print(f"[KG] Stored graph for {nomor}: {len(data.get('entitas',[]))} entities, {len(data.get('relasi',[]))} edges")
    except Exception as e:
        print(f"[KG] DB write failed for {nomor}: {e}")

import auth
import history
import internal_docs
import templates

app.include_router(auth.router, prefix="/api/auth", tags=["auth"])
app.include_router(history.router, prefix="/api", tags=["history"])
app.include_router(internal_docs.router, prefix="/api", tags=["internal_docs"])
app.include_router(templates.router, prefix="/api", tags=["templates"])

# Serve local PDFs as static files
PDFS_DIR = os.path.join(BASE_DIR, "data", "pdfs")
os.makedirs(PDFS_DIR, exist_ok=True)
# NOTE: We do NOT use app.mount(StaticFiles) because Starlette mounts bypass CORS middleware.
# Instead we use a regular route below (/api/pdf/{filename}) which correctly gets CORS headers.

# Setup CORS to allow frontend communication
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # For prototype purposes
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Models
class ChatRequest(BaseModel):
    messages: List[dict]  # [{"role": "user", "content": "..."}, ...]

class Source(BaseModel):
    id: str
    jenis: str
    nomor: str
    sektor: str
    judul: str
    snippet: str

class ChatResponse(BaseModel):
    answer: str
    sources: List[Source]

class AnalyzeRequest(BaseModel):
    reg_id: Optional[str] = None
    nomor: str
    judul: str
    filename: Optional[str] = None

    @field_validator('reg_id', 'nomor', 'judul', mode='before')
    @classmethod
    def coerce_to_str(cls, v: Any) -> Optional[str]:
        if v is None:
            return None
        return str(v)

class ConfirmPendingRequest(BaseModel):
    klasifikasi: str

@app.post("/api/repository/pending/{doc_id}/confirm")
async def confirm_pending_document(
    doc_id: str,
    req: ConfirmPendingRequest,
    background_tasks: BackgroundTasks,
    current_user: dict = Depends(auth.require_role("sekretaris perusahaan"))
):
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    c = conn.cursor()
    
    c.execute("SELECT local_path, nomor, judul, jenis, sektor FROM regulations WHERE id = ? AND status = 'Menunggu Konfirmasi'", (doc_id,))
    doc = c.fetchone()
    
    if not doc:
        conn.close()
        raise HTTPException(status_code=404, detail="Dokumen pending tidak ditemukan.")
        
    local_path, nomor, judul, jenis, sektor = doc
    
    # Update DB
    c.execute("UPDATE regulations SET status = 'Berlaku', klasifikasi = ? WHERE id = ?", (req.klasifikasi, doc_id))
    conn.commit()
    conn.close()
    
    # Trigger vector DB injection and KG asynchronously
    filename = os.path.basename(local_path)
    background_tasks.add_task(ingest_document_background, local_path, doc_id, filename, nomor, jenis, sektor, 'Berlaku', req.klasifikasi)
    
    return {"message": "Dokumen berhasil dikonfirmasi dan dimasukkan ke repositori."}

class AccessGrantRequest(BaseModel):
    granted_to: str
    reason: str
    expires_at: Optional[str] = None

# ChromaDB Cache
_collection = None
_glm_session: Optional[requests.Session] = None

GENERIC_GLM_ERROR_MESSAGE = (
    "Layanan AI sedang mengalami gangguan koneksi sementara. "
    "Silakan coba lagi dalam beberapa saat."
)

def get_chroma_collection():
    global _collection
    if _collection is None:
        client = chromadb.PersistentClient(path=CHROMA_DB_DIR)
        _collection = client.get_or_create_collection(
            name="ojk_regulations",
            metadata={"hnsw:space": "cosine"}
        )
    return _collection

def get_glm_session() -> requests.Session:
    """Build a reusable HTTP session with transport-level retries."""
    global _glm_session
    if _glm_session is None:
        session = requests.Session()
        import urllib3
        urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
        session.verify = False
        
        retry_cfg = Retry(
            total=2,
            connect=2,
            read=2,
            status=2,
            backoff_factor=0.6,
            status_forcelist=[429, 500, 502, 503, 504],
            allowed_methods=frozenset(["POST"]),
            raise_on_status=False,
        )
        adapter = HTTPAdapter(max_retries=retry_cfg, pool_connections=10, pool_maxsize=20)
        session.mount("https://", adapter)
        session.mount("http://", adapter)
        _glm_session = session
    return _glm_session

def is_transient_network_error(err: Exception | None) -> bool:
    if err is None:
        return False
    text = str(err).lower()
    transient_markers = [
        "name resolution",
        "failed to resolve",
        "getaddrinfo failed",
        "remote end closed connection",
        "connection aborted",
        "connection reset",
        "temporarily unavailable",
        "timed out",
        "timeout",
    ]
    return any(marker in text for marker in transient_markers)

def extract_pasal_items(text: str, max_items: int = 20, max_chars_per_pasal: int = 900) -> List[dict]:
    """
    Extract Pasal blocks from document text.
    Supports both Arabic and Roman numerals, e.g., "Pasal 3" and "PASAL III".
    """
    # Normalize line endings and spacing around line breaks for stable regex matching.
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    normalized = re.sub(r'[ \t]+\n', '\n', normalized)
    normalized = re.sub(r'\n{3,}', '\n\n', normalized)

    # Strict heading match first (line-based)
    pasal_matches = list(re.finditer(r'(?im)^\s*pasal\s+((?:\d+|[ivxlcdm]+))\b', normalized))
    # Fallback for documents where line breaks are collapsed by extractor
    if not pasal_matches:
        pasal_matches = list(re.finditer(r'(?i)\bpasal\s+((?:\d+|[ivxlcdm]+))\b', normalized))

    pasal_items = []
    for idx, match in enumerate(pasal_matches):
        pasal_no_raw = match.group(1).upper()
        pasal_label = f"Pasal {pasal_no_raw}"
def extract_pasal_items(text: str, max_items: int = 20, max_chars_per_pasal: int = 900) -> List[dict]:
    """
    Primary: regex-based Pasal extractor (fast, zero API cost).
    Used as the first attempt. LLM-based extractor is used as fallback in check_compliance.
    """
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    normalized = re.sub(r'[ \t]+\n', '\n', normalized)
    normalized = re.sub(r'\n{3,}', '\n\n', normalized)

    pasal_matches = list(re.finditer(r'(?im)^\s*pasal\s+((?:\d+|[ivxlcdm]+))\b', normalized))
    if not pasal_matches:
        pasal_matches = list(re.finditer(r'(?i)\bpasal\s+((?:\d+|[ivxlcdm]+))\b', normalized))

    pasal_items = []
    for idx, match in enumerate(pasal_matches):
        pasal_no_raw = match.group(1).upper()
        pasal_label = f"Pasal {pasal_no_raw}"
        start_idx = match.start()
        end_idx = pasal_matches[idx + 1].start() if idx + 1 < len(pasal_matches) else len(normalized)
        pasal_text = normalized[start_idx:end_idx].strip()
        if len(pasal_text) >= 20:
            pasal_items.append({
                "pasal": pasal_label,
                "deskripsi": pasal_text[:max_chars_per_pasal]
            })
        if len(pasal_items) >= max_items:
            break

    return pasal_items


def llm_extract_pasal_items(full_text: str, max_items: int = 20) -> List[dict]:
    """
    Fallback LLM-based clause extractor for documents where regex fails
    (bilingual contracts, non-standard headings, Article X format, etc.).
    """
    snippet = full_text[:8000]
    prompt = (
        "Dokumen hukum berikut adalah perjanjian/kontrak. "
        "Ekstrak SETIAP klausul atau pasal sebagai daftar JSON. "
        "Format: [{\"pasal\": \"Pasal 1\", \"deskripsi\": \"teks lengkap klausul...\"}]. "
        "Jika tidak ada heading Pasal, gunakan nomor urut klausul. "
        "Kembalikan HANYA JSON array, tanpa teks lain.\n\n"
        f"DOKUMEN:\n{snippet}"
    )
    try:
        raw = call_glm([{"role": "user", "content": prompt}], temperature=0.0, timeout=60)
        raw = re.sub(r'```json|```', '', raw).strip()
        items = json.loads(raw)
        if isinstance(items, list):
            return [{"pasal": str(i.get("pasal", f"Klausul {n+1}")),
                     "deskripsi": str(i.get("deskripsi", ""))[:900]}
                    for n, i in enumerate(items[:max_items])
                    if i.get("deskripsi")]
    except Exception as e:
        print(f"[LLM Pasal Extractor] Failed: {e}")
    return []


def _regex_extract_date_duration(text: str) -> dict:
    """
    Regex + keyword fallback to extract document creation date and duration
    when the LLM misses them. Handles common Indonesian legal phrasings.
    """
    import re as _re
    result = {"tanggal_pembuatan": None, "durasi_perjanjian_bulan": None}

    # ── 1. Written-number month map ─────────────────────────────────────────
    BULAN_MAP = {
        "januari": 1, "februari": 2, "maret": 3, "april": 4,
        "mei": 5, "juni": 6, "juli": 7, "agustus": 8,
        "september": 9, "oktober": 10, "november": 11, "desember": 12,
    }
    ANGKA_MAP = {
        "satu": 1, "dua": 2, "tiga": 3, "empat": 4, "lima": 5,
        "enam": 6, "tujuh": 7, "delapan": 8, "sembilan": 9, "sepuluh": 10,
        "sebelas": 11, "dua belas": 12, "dua puluh empat": 24,
        "tiga puluh enam": 36, "empat puluh delapan": 48, "enam puluh": 60,
    }

    lowered = text.lower()

    # ── 2. Creation date: "hari ini [weekday] tanggal [X] [bulan] [tahun]" ──
    # Handles both digit and written day
    date_patterns = [
        # "tanggal delapan belas bulan februari tahun dua ribu dua puluh satu"
        r"tanggal\s+([\w\s]+?)\s+bulan\s+(januari|februari|maret|april|mei|juni|juli|agustus|september|oktober|november|desember)\s+tahun\s+([\w\s]+)",
        # "tanggal 18 februari 2021"
        r"tanggal\s+(\d{1,2})\s+(januari|februari|maret|april|mei|juni|juli|agustus|september|oktober|november|desember)\s+(\d{4})",
        # "dibuat di ... pada tanggal 18-02-2021"
        r"(?:dibuat|ditandatangani|berlaku)\s+(?:pada\s+)?tanggal\s+(\d{1,2})[/\-\.](\d{1,2})[/\-\.](\d{4})",
        # "Jakarta, 18 Februari 2021"
        r"(?:jakarta|bandung|surabaya|medan|makassar|semarang|depok|bogor|bekasi|tangerang)[,\s]+(\d{1,2})\s+(januari|februari|maret|april|mei|juni|juli|agustus|september|oktober|november|desember)\s+(\d{4})",
    ]

    WRITTEN_NUMS_DAY = {
        "satu": 1, "dua": 2, "tiga": 3, "empat": 4, "lima": 5,
        "enam": 6, "tujuh": 7, "delapan": 8, "sembilan": 9, "sepuluh": 10,
        "sebelas": 11, "dua belas": 12, "tiga belas": 13, "empat belas": 14,
        "lima belas": 15, "enam belas": 16, "tujuh belas": 17, "delapan belas": 18,
        "sembilan belas": 19, "dua puluh": 20, "dua puluh satu": 21,
        "dua puluh dua": 22, "dua puluh tiga": 23, "dua puluh empat": 24,
        "dua puluh lima": 25, "dua puluh enam": 26, "dua puluh tujuh": 27,
        "dua puluh delapan": 28, "dua puluh sembilan": 29, "tiga puluh": 30,
        "tiga puluh satu": 31,
    }
    WRITTEN_YEARS = {
        "dua ribu dua puluh": 2020, "dua ribu dua puluh satu": 2021,
        "dua ribu dua puluh dua": 2022, "dua ribu dua puluh tiga": 2023,
        "dua ribu dua puluh empat": 2024, "dua ribu dua puluh lima": 2025,
        "dua ribu dua puluh enam": 2026, "dua ribu sembilan belas": 2019,
        "dua ribu delapan belas": 2018, "dua ribu tujuh belas": 2017,
    }

    for pat in date_patterns:
        m = _re.search(pat, lowered)
        if m:
            groups = m.groups()
            try:
                if len(groups) == 3:
                    g0, g1, g2 = [g.strip() for g in groups]
                    # Resolve day
                    day = int(g0) if g0.isdigit() else WRITTEN_NUMS_DAY.get(g0, None)
                    # Resolve month
                    month_str = g1.lower()
                    month = BULAN_MAP.get(month_str, None)
                    if month is None and g1.isdigit():
                        month = int(g1)
                    # Resolve year
                    year = int(g2) if g2.isdigit() else WRITTEN_YEARS.get(g2.strip(), None)
                    if day and month and year and 2000 <= year <= 2035:
                        from datetime import date
                        result["tanggal_pembuatan"] = date(year, month, day).strftime("%Y-%m-%d")
                        break
            except Exception:
                continue

    # ── 3. Duration: look inside JANGKA WAKTU section first, then globally ──
    # First, try to extract just the text of the "Jangka Waktu" clause/section
    jw_section = ""
    jw_match = _re.search(
        r'(?:pasal\s*\d+\s*)?jangka\s+waktu[\w\s]*?\n(.{0,800})',
        lowered, _re.DOTALL
    )
    if jw_match:
        jw_section = jw_match.group(0)

    search_texts = [jw_section, lowered] if jw_section else [lowered]

    dur_patterns = [
        # "minimal selama 3 (tiga) tahun"  ← from your PKS example
        r"(?:minimal|paling\s+sedikit)?\s*selama\s+(\d+)\s*(?:\([\w\s]+\))?\s*(tahun|bulan)",
        # "adalah minimal selama 3 (tiga) tahun"
        r"adalah\s+(?:minimal\s+)?selama\s+(\d+)\s*(?:\([\w\s]+\))?\s*(tahun|bulan)",
        # "berlaku selama 2 (dua) tahun"
        r"(?:berlaku|berlangsung)\s+selama\s+(\d+)\s*(?:\([\w\s]+\))?\s*(tahun|bulan)",
        # "jangka waktu ... adalah/selama 12 (dua belas) bulan"
        r"jangka\s+waktu\s+[\w\s]{0,40}?(?:adalah|selama|yaitu|:)?\s*(\d+)\s*(?:\([\w\s]+\))?\s*(tahun|bulan)",
        # "selama dua tahun" (written number)
        r"selama\s+(satu|dua|tiga|empat|lima|enam|tujuh|delapan|sembilan|sepuluh|sebelas|dua belas|dua puluh empat|tiga puluh enam)\s*(tahun|bulan)",
        # "masa berlaku ... 1 (satu) tahun"
        r"masa\s+berlaku\s+[\w\s,]{0,60}?(\d+)\s*(?:\([\w\s]+\))?\s*(tahun|bulan)",
        # "perjanjian ini berlangsung selama 24 bulan"
        r"perjanjian\s+ini\s+(?:akan\s+)?berlangsung\s+selama\s+(\d+)\s*(?:\([\w\s]+\))?\s*(tahun|bulan)",
    ]

    for search_text in search_texts:
        if result["durasi_perjanjian_bulan"]:
            break
        for pat in dur_patterns:
            m = _re.search(pat, search_text)
            if m:
                val_str, unit = m.group(1).strip(), m.group(2).strip()
                try:
                    val = int(val_str) if val_str.isdigit() else ANGKA_MAP.get(val_str, None)
                    if val:
                        result["durasi_perjanjian_bulan"] = val * 12 if unit == "tahun" else val
                        break
                except Exception:
                    continue

    return result


def understand_document(text: str) -> dict:
    """
    Pass 0A — Document Metadata Understanding.
    Focused ONLY on: doc type, parties, sector, subject.
    Date and duration are extracted by dedicated functions below.
    Uses only the preamble (first 6000 chars) for speed and accuracy.
    """
    snippet = text[:6000]
    prompt = (
        "Baca bagian awal dokumen hukum Indonesia ini dan identifikasi informasi berikut. "
        "Kembalikan HANYA JSON tanpa teks tambahan:\n"
        '{\n'
        '  "jenis_dokumen": "PKS / NDA / Perjanjian Kerja / Kontrak Layanan / dll",\n'
        '  "pihak_pertama": "nama lengkap perusahaan/entitas pihak pertama",\n'
        '  "pihak_kedua": "nama lengkap perusahaan/entitas pihak kedua",\n'
        '  "sektor_bisnis": "teknologi / keuangan / ketenagakerjaan / perbankan / dll",\n'
        '  "pokok_perjanjian": "deskripsi singkat isi perjanjian dalam 1 kalimat"\n'
        '}\n\n'
        f"DOKUMEN (BAGIAN AWAL):\n{snippet}"
    )
    try:
        raw = call_glm([{"role": "user", "content": prompt}], temperature=0.0, timeout=30)
        raw = re.sub(r'```json|```', '', raw).strip()
        result = json.loads(raw)
        print(f"[DocContext] Metadata: {result}")
        return result
    except Exception as e:
        print(f"[DocContext] Metadata LLM failed: {e}")
        return {
            "jenis_dokumen": "Kontrak",
            "pihak_pertama": "Pihak Pertama",
            "pihak_kedua": "Pihak Kedua",
            "sektor_bisnis": "umum",
            "pokok_perjanjian": "tidak teridentifikasi",
        }


def extract_signing_date(text: str) -> str | None:
    """
    Pass 0B — Dedicated signing date extraction.
    Strategy:
      1. Scan preamble (first 3000 chars) where opening date is stated
      2. Scan signature block (last 2000 chars) where city+date appears
      3. Regex fallback across both windows
    """
    preamble   = text[:3000]
    sig_block  = text[-2000:]
    combined   = preamble + "\n\n[...BAGIAN AKHIR DOKUMEN...]\n\n" + sig_block

    prompt = (
        "Dari teks dokumen hukum Indonesia berikut (bagian AWAL dan AKHIR dokumen), "
        "temukan tanggal penandatanganan atau pembuatan perjanjian ini.\n"
        "Tanggal ini biasanya muncul dalam bentuk:\n"
        "- 'dibuat/ditandatangani pada hari ... tanggal [X] bulan [Y] tahun [Z]'\n"
        "- '[Kota], [tanggal] [bulan] [tahun]' (contoh: 'Jakarta, 18 Februari 2021')\n"
        "- 'tanggal delapan belas bulan februari tahun dua ribu dua puluh satu'\n"
        "Angka boleh berupa kata (delapan belas = 18, dua ribu dua puluh satu = 2021).\n"
        "Kembalikan HANYA JSON: {\"tanggal\": \"YYYY-MM-DD\"} atau {\"tanggal\": null} jika tidak ditemukan.\n\n"
        f"TEKS:\n{combined}"
    )
    llm_date = None
    try:
        raw = call_glm([{"role": "user", "content": prompt}], temperature=0.0, timeout=25)
        raw = re.sub(r'```json|```', '', raw).strip()
        parsed = json.loads(raw)
        llm_date = parsed.get("tanggal")
        if llm_date:
            # Validate format
            datetime.strptime(llm_date, "%Y-%m-%d")
            print(f"[SigningDate] LLM found: {llm_date}")
    except Exception as e:
        print(f"[SigningDate] LLM failed or invalid date: {e}")
        llm_date = None

    if llm_date:
        return llm_date

    # Regex fallback on preamble + sig block
    regex_result = _regex_extract_date_duration(preamble + sig_block)
    if regex_result.get("tanggal_pembuatan"):
        print(f"[SigningDate] Regex fallback: {regex_result['tanggal_pembuatan']}")
        return regex_result["tanggal_pembuatan"]

    return None


def extract_duration_months(text: str) -> int | None:
    """
    Pass 0C — Dedicated duration extraction.
    Strategy:
      1. Isolate the 'JANGKA WAKTU' section using regex on section headers
      2. Feed ONLY that section (+ small buffer) to the LLM
      3. If no section found, fall back to scanning full text with LLM
      4. Regex fallback on the JANGKA WAKTU section or full text
    """
    lowered = text.lower()

    # ── Step 1: Find and isolate the JANGKA WAKTU section ───────────────────
    # Matches: PASAL N\nJANGKA WAKTU... or just JANGKA WAKTU PELAKSANAAN...
    jw_pattern = re.compile(
        r'(?:pasal\s*\d+\s*[\n\r]+)?'
        r'jangka\s+waktu[\w\s]*?[\n\r]'
        r'(.{100,1500}?)'
        r'(?=pasal\s*\d+|\Z)',
        re.IGNORECASE | re.DOTALL
    )
    jw_match = jw_pattern.search(text)
    jw_section = jw_match.group(0) if jw_match else ""

    if jw_section:
        print(f"[Duration] Found JANGKA WAKTU section ({len(jw_section)} chars)")
        context_for_llm = jw_section[:1500]
    else:
        print("[Duration] No JANGKA WAKTU section found, using full text")
        context_for_llm = text[:20000]

    prompt = (
        "Dari teks berikut (diambil dari klausul JANGKA WAKTU perjanjian), "
        "temukan durasi/jangka waktu berlakunya perjanjian ini.\n"
        "Durasi biasanya dinyatakan dalam bentuk:\n"
        "- 'berlaku selama X tahun/bulan'\n"
        "- 'minimal selama X (Y) tahun'\n"
        "- 'jangka waktu ... adalah X bulan'\n"
        "- 'masa berlaku X tahun'\n"
        "PENTING: Konversi semua ke BULAN (1 tahun = 12 bulan, 2 tahun = 24, 3 tahun = 36).\n"
        "Kembalikan HANYA JSON: {\"durasi_bulan\": <integer>} atau {\"durasi_bulan\": null} jika tidak ditemukan.\n\n"
        f"TEKS:\n{context_for_llm}"
    )
    llm_duration = None
    try:
        raw = call_glm([{"role": "user", "content": prompt}], temperature=0.0, timeout=25)
        raw = re.sub(r'```json|```', '', raw).strip()
        parsed = json.loads(raw)
        val = parsed.get("durasi_bulan")
        if val and isinstance(val, (int, float)) and 1 <= int(val) <= 600:
            llm_duration = int(val)
            print(f"[Duration] LLM found: {llm_duration} bulan")
    except Exception as e:
        print(f"[Duration] LLM failed: {e}")
        llm_duration = None

    if llm_duration:
        return llm_duration

    # Regex fallback — try JANGKA WAKTU section first, then full text
    for search_src in ([jw_section, text] if jw_section else [text]):
        regex_result = _regex_extract_date_duration(search_src[:30000])
        if regex_result.get("durasi_perjanjian_bulan"):
            dur = regex_result["durasi_perjanjian_bulan"]
            print(f"[Duration] Regex fallback: {dur} bulan")
            return dur

    return None

def extract_explicit_end_date(text: str) -> str | None:
    """
    Pass 0D — Extracts an explicit end date from the JANGKA WAKTU section.
    Used as a fallback when 'start_date + duration_months' fails.
    Looks for 'sampai dengan tanggal X', 'berakhir pada Y'.
    """
    jw_pattern = re.compile(
        r'(?:pasal\s*\d+\s*[\n\r]+)?'
        r'jangka\s+waktu[\w\s]*?[\n\r]'
        r'(.{100,1500}?)'
        r'(?=pasal\s*\d+|\Z)',
        re.IGNORECASE | re.DOTALL
    )
    jw_match = jw_pattern.search(text)
    context_for_llm = jw_match.group(0)[:1500] if jw_match else text[:20000]

    prompt = (
        "Dari teks berikut, temukan TANGGAL BERAKHIR (kedaluwarsa) perjanjian secara eksplisit.\n"
        "Cari frasa seperti:\n"
        "- 'berlaku ... sampai dengan tanggal [X] bulan [Y] tahun [Z]'\n"
        "- 'berakhir pada tanggal [tanggal]'\n"
        "Konversi ke format YYYY-MM-DD. Angka boleh berupa huruf (dua ribu dua puluh empat = 2024).\n"
        "Kembalikan HANYA JSON: {\"tanggal_berakhir\": \"YYYY-MM-DD\"} atau {\"tanggal_berakhir\": null} jika tidak disebutkan secara spesifik.\n\n"
        f"TEKS:\n{context_for_llm}"
    )
    try:
        raw = call_glm([{"role": "user", "content": prompt}], temperature=0.0, timeout=25)
        raw = re.sub(r'```json|```', '', raw).strip()
        parsed = json.loads(raw)
        end_date = parsed.get("tanggal_berakhir")
        if end_date:
            datetime.strptime(end_date, "%Y-%m-%d")
            print(f"[ExplicitEndDate] LLM found: {end_date}")
            return end_date
    except Exception as e:
        print(f"[ExplicitEndDate] LLM failed: {e}")
        
    # Regex fallback for explicit end dates
    lowered = context_for_llm.lower()
    end_date_patterns = [
        r"(?:sampai\s+dengan|berakhir\s+pada)(?:\s+tanggal)?\s+(\d{1,2})\s+(januari|februari|maret|april|mei|juni|juli|agustus|september|oktober|november|desember)\s+(\d{4})",
        r"(?:sampai\s+dengan|berakhir\s+pada)(?:\s+tanggal)?\s+(\d{1,2})[/\-\.](\d{1,2})[/\-\.](\d{4})"
    ]
    
    BULAN_MAP = {
        "januari": 1, "februari": 2, "maret": 3, "april": 4,
        "mei": 5, "juni": 6, "juli": 7, "agustus": 8,
        "september": 9, "oktober": 10, "november": 11, "desember": 12,
    }
    
    for pat in end_date_patterns:
        m = re.search(pat, lowered)
        if m:
            groups = m.groups()
            try:
                if len(groups) == 3:
                    day = int(groups[0])
                    month_str = groups[1].lower()
                    month = BULAN_MAP.get(month_str) if not month_str.isdigit() else int(month_str)
                    year = int(groups[2])
                    if day and month and year and 2000 <= year <= 2050:
                        from datetime import date
                        res = date(year, month, day).strftime("%Y-%m-%d")
                        print(f"[ExplicitEndDate] Regex fallback found: {res}")
                        return res
            except Exception:
                continue

    return None



def is_regulation_relevant(pasal_text: str, reg_text: str, reg_name: str) -> bool:
    """
    Relevance confirmation micro-call.
    Fast Yes/No check: does this regulation actually apply to this clause?
    Prevents wrong-domain regulations from being cited confidently.
    """
    prompt = (
        f"Apakah regulasi '{reg_name}' berikut SECARA LANGSUNG relevan untuk mengevaluasi "
        f"klausul kontrak berikut? Jawab hanya 'YA' atau 'TIDAK'.\n\n"
        f"KLAUSUL:\n{pasal_text[:400]}\n\n"
        f"REGULASI:\n{reg_text[:600]}"
    )
    try:
        ans = call_glm([{"role": "user", "content": prompt}], temperature=0.0, timeout=20)
        return "YA" in ans.upper()
    except Exception:
        return True  # default: assume relevant if check fails

def vlm_extract_page_image(page) -> str:
    """
    Render a PyMuPDF page as a PNG image and return it as a base64-encoded string.
    Uses 288 DPI (3x scale) to ensure clear distinction of characters like 'j' and 'l'.
    """
    import fitz
    mat = fitz.Matrix(3, 3)  # 3x = ~288 DPI
    pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
    return base64.b64encode(pix.tobytes("png")).decode("utf-8")


def call_glm_vision(image_b64: str, prompt: str, timeout: int = 90) -> str:
    """
    Send a page image to llama-4-maverick via the vision (multimodal) API.
    Returns the extracted text content.
    """
    if API_BASE_URL.endswith("/chat/completions"):
        url = API_BASE_URL
    else:
        url = f"{API_BASE_URL}/chat/completions"

    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": VLM_MODEL,
        "messages": [{
            "role": "user",
            "content": [
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{image_b64}"}
                },
                {
                    "type": "text",
                    "text": prompt
                }
            ]
        }],
        "temperature": 0.0,
        "stream": False
    }
    resp = get_glm_session().post(url, json=payload, headers=headers, timeout=(15, timeout))
    if resp.ok:
        return resp.json()["choices"][0]["message"]["content"]
    print(f"VLM page extraction failed: {resp.status_code} {resp.text[:200]}")
    return ""


VLM_PAGE_PROMPT = (
    "Ekstrak SEMUA teks dari dokumen ini secara akurat. "
    "Perbaiki kesalahan ejaan visual/typo hasil OCR (misalnya huruf 'l' yang seharusnya 'J') "
    "agar membentuk kalimat bahasa Indonesia yang baku dan masuk akal. "
    "Pertahankan struktur asli seperti nomor pasal dan ayat. "
    "Jangan tambahkan penjelasan atau komentar di luar teks dokumen."
)



def is_text_garbled(text: str) -> bool:
    """
    Detect font-encoding corruption in PyMuPDF-extracted text.
    Returns True if the text looks garbled/unreadable.

    Three independent signals — triggering ANY ONE marks the page as corrupted:
    1. Words with ZERO vowels (length >= 3): e.g. 'KBWJBN', 'PHRBN'
       These are statistically impossible in Indonesian/English legal text.
    2. Consonant clusters of 3+ in >= 8% of words: e.g. 'KBWAJIBAN' (K-B-W = 3)
       Previous threshold was 4, which missed these.
    3. Global vowel ratio < 20%: whole-page signal for systematic encoding failure.
    """
    if len(text) < 30:
        return False

    words = text.split()
    if not words:
        return False

    VOWELS    = set("aeiouAEIOU")
    CONSONANTS = set("bcdfghjklmnpqrstvwxyzBCDFGHJKLMNPQRSTVWXYZ")

    cluster_words  = 0   # words with 3+ consecutive consonants
    no_vowel_words = 0   # words with NO vowels at all (len >= 3)

    for w in words:
        alpha = [c for c in w if c.isalpha()]
        if not alpha:
            continue

        # Signal 1: zero-vowel words (strong indicator of garbling)
        if len(alpha) >= 3 and not any(c in VOWELS for c in alpha):
            no_vowel_words += 1

        # Signal 2: consonant cluster (3+ threshold, down from 4)
        if len(alpha) >= 3:
            run = max_run = 0
            for c in alpha:
                run = run + 1 if c in CONSONANTS else 0
                max_run = max(max_run, run)
            if max_run >= 3:
                cluster_words += 1

    total_words    = max(len(words), 1)
    cluster_ratio  = cluster_words  / total_words
    no_vowel_ratio = no_vowel_words / total_words

    # Signal 3: global vowel ratio
    alpha_chars = [c for c in text if c.isalpha()]
    vowel_ratio = (sum(1 for c in alpha_chars if c in VOWELS) / len(alpha_chars)
                   if alpha_chars else 0.0)

    reasons = []
    if no_vowel_ratio > 0.05:  reasons.append(f"no-vowel-words={no_vowel_ratio:.2f}")
    if cluster_ratio  > 0.08:  reasons.append(f"cluster_ratio={cluster_ratio:.2f}")
    if vowel_ratio    < 0.20:  reasons.append(f"vowel_ratio={vowel_ratio:.2f}")

    garbled = bool(reasons)
    if garbled:
        print(f"  [QC] Garbled text detected: {', '.join(reasons)}")
    return garbled


def is_text_readable_llm(text: str) -> bool:
    """
    Fast LLM check to see if PyMuPDF text is readable Indonesian or garbled.
    """
    if len(text) < 50:
        return True # Too short, assume readable or handled by length check
    
    snippet = text[:500]
    prompt = (
        "Apakah teks berikut merupakan teks bahasa Indonesia/Inggris yang dapat dibaca, "
        "ataukah teks tersebut rusak/acak (garbled) karena kesalahan font encoding? "
        "Jawab HANYA dengan 'BISA DIBACA' atau 'RUSAK'.\n\n"
        f"TEKS:\n{snippet}"
    )
    try:
        ans = call_glm([{"role": "user", "content": prompt}], temperature=0.0, timeout=20)
        return "BISA DIBACA" in ans.upper()
    except Exception:
        return True # Default to readable if check fails


def extract_text_hybrid(pdf_path: str, digital_threshold: int = 80,
                        force_vlm: bool = False) -> str:
    """
    Hybrid PDF text extractor:
    - force_vlm=True : ALWAYS sends every page through VLM (used by compliance checker
                       for maximum accuracy regardless of PDF type).
    - force_vlm=False: Smart routing — clean digital pages use PyMuPDF, scanned/
                       font-corrupted pages use VLM (used by ingestion pipeline).
    """
    import fitz
    doc = fitz.open(pdf_path)
    full_text_parts = []
    vlm_pages = 0
    digital_pages = 0

    # Fast-pass LLM check on the first few pages to detect global font corruption
    global_garbled = False
    if not force_vlm:
        for p in range(min(3, len(doc))):
            pt = doc.load_page(p).get_text("text").strip()
            if len(pt) > digital_threshold:
                if is_text_garbled(pt) or not is_text_readable_llm(pt):
                    global_garbled = True
                break
        if global_garbled:
            print("  [QC] Global font corruption detected. Routing entire document to VLM.")

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        digital_text = page.get_text("text").strip()

        use_vlm = False
        reason = ""

        if force_vlm:
            use_vlm = True
            reason = "force_vlm=True (compliance mode)"
        elif global_garbled:
            use_vlm = True
            reason = "global font corruption"
        elif len(digital_text) < digital_threshold:
            use_vlm = True
            reason = f"scanned/empty ({len(digital_text)} chars)"
        elif is_text_garbled(digital_text):
            use_vlm = True
            reason = "font encoding corruption detected"

        if not use_vlm:
            cleaned = re.sub(r'([^\n])\n([^\n])', r'\1 \2', digital_text)
            full_text_parts.append(cleaned)
            digital_pages += 1
        else:
            print(f"  [VLM] Page {page_num + 1} -> {reason}")
            try:
                img_b64 = vlm_extract_page_image(page)
                vlm_text = call_glm_vision(img_b64, VLM_PAGE_PROMPT, timeout=60)
                if vlm_text.strip():
                    full_text_parts.append(vlm_text.strip())
                    vlm_pages += 1
                else:
                    print(f"  [VLM] Page {page_num + 1} returned empty — keeping original.")
                    full_text_parts.append(digital_text)
            except Exception as e:
                print(f"  [VLM] Page {page_num + 1} vision error: {e} — keeping original.")
                full_text_parts.append(digital_text)

    doc.close()
    print(f"  [Hybrid] Done: {digital_pages} direct + {vlm_pages} VLM pages")
    return "\n\n".join(full_text_parts)


def call_glm(messages: list, temperature: float = 0.1, timeout: int = 90) -> str:
    """
    Unified GLM API caller (OpenAI-compatible).
    Returns the assistant message content string.
    Raises HTTPException on failure.
    """
    # Handle cases where user might have included /chat/completions in the base URL
    if API_BASE_URL.endswith("/chat/completions"):
        url = API_BASE_URL
    else:
        url = f"{API_BASE_URL}/chat/completions"
    
    # DEBUG: See what the server is actually using
    print(f"DEBUG: Calling LLM at URL: {url}")
    print(f"DEBUG: Model: {GLM_MODEL}")

    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": GLM_MODEL,
        "messages": messages,
        "temperature": temperature,
        "stream": False
    }
    max_attempts = max(1, GLM_MAX_ATTEMPTS)
    last_error: Exception | None = None
    session = get_glm_session()

    for attempt in range(1, max_attempts + 1):
        try:
            # Tuple timeout: (connect_timeout, read_timeout)
            resp = session.post(url, json=payload, headers=headers, timeout=(15, timeout))

            if not resp.ok:
                # Include truncated upstream body to speed up debugging bad credentials/model/payload.
                body_snippet = (resp.text or "")[:300]
                print(f"GLM non-2xx response: status={resp.status_code}, body={body_snippet}")
                raise HTTPException(
                    status_code=500,
                    detail=GENERIC_GLM_ERROR_MESSAGE
                )

            return resp.json()["choices"][0]["message"]["content"]
        except HTTPException:
            raise
        except requests.exceptions.RequestException as e:
            last_error = e
            print(f"GLM request attempt {attempt}/{max_attempts} failed: {e}")
            if attempt < max_attempts:
                # Exponential backoff to absorb transient upstream disconnects.
                sleep_s = GLM_RETRY_BACKOFF_BASE * (2 ** (attempt - 1))
                time.sleep(sleep_s)
                continue
            break
        except (KeyError, IndexError, TypeError) as e:
            print(f"Unexpected GLM response format: {e}")
            raise HTTPException(status_code=500, detail=GENERIC_GLM_ERROR_MESSAGE)

    if is_transient_network_error(last_error):
        raise HTTPException(status_code=503, detail=GENERIC_GLM_ERROR_MESSAGE)
    raise HTTPException(status_code=500, detail=GENERIC_GLM_ERROR_MESSAGE)


def retrieve_contexts(query: str, current_user: dict, n_results=5, doc_category=None, hybrid_cag=True) -> List[dict]:
    """Searches Vector DB and returns structured raw dictionaries.
       If hybrid_cag=True, uses RAG to find the best document, then CAG to inject the full text.
       Enforces FR-17 by strictly filtering out classified documents the user cannot access.
    """
    collection = get_chroma_collection()
    user_id = current_user.get("id")
    role_level = auth.get_role_level(current_user.get("role", "pengguna"))
    
    # 1. Build Base Allowed Classifications
    allowed_klasifikasi = ["Umum"]
    if role_level >= 2:
        allowed_klasifikasi.append("Rahasia")
    if role_level >= 3:
        allowed_klasifikasi.append("Terbatas")
        
    # 2. Fetch Explicit Access Grants & Document Classification Map
    granted_doc_ids = set()
    doc_klasifikasi_map = {}
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    try:
        conn = sqlite3.connect(db_path)
        c = conn.cursor()
        c.execute("SELECT doc_id FROM access_grants WHERE granted_to = ? AND expires_at >= datetime('now')", (user_id,))
        for row in c.fetchall():
            granted_doc_ids.add(row[0])
            
        c.execute("SELECT id, klasifikasi FROM regulations")
        for row in c.fetchall():
            doc_klasifikasi_map[row[0]] = row[1] or "Umum"
        conn.close()
    except Exception as e:
        print(f"Error fetching access grants for context retrieval: {e}")

    def is_allowed(reg_id: str) -> bool:
        if not reg_id:
            return True # Allow edge case for very old unstructured data
        if reg_id in granted_doc_ids:
            return True
        doc_klas = doc_klasifikasi_map.get(reg_id, "Umum")
        return doc_klas in allowed_klasifikasi

    # Over-fetch for Post-Retrieval Filtering
    overfetch_n = 50
    where_clause = {"$or": [{"visibility": "public"}, {"user_id": user_id}]}
    if doc_category:
        where_clause = {
            "$and": [
                where_clause,
                {"doc_category": {"$in": [doc_category, "UMUM"]}}
            ]
        }
        
    contexts = []
    
    if hybrid_cag:
        # 1. RAG Discovery: Find the single most relevant chunk
        discovery = collection.query(query_texts=[query], n_results=overfetch_n, where=where_clause)
        
        best_reg_id = None
        best_top_meta = None
        
        if discovery and discovery['documents'] and len(discovery['documents'][0]) > 0:
            # Post-Filter to find the top AUTHORIZED document
            for idx in range(len(discovery['documents'][0])):
                meta = discovery['metadatas'][0][idx]
                reg_id = meta.get("reg_id")
                
                if is_allowed(reg_id):
                    best_reg_id = reg_id
                    best_top_meta = meta
                    break # Found the highest ranked authorized document
            
            if best_reg_id:
                # 2. CAG Injection: Fetch all chunks for this specific document
                doc_results = collection.get(where={"reg_id": best_reg_id}, include=["documents"])
                
                if doc_results and doc_results['documents']:
                    # Reconstruct full text
                    full_text = "\n".join(doc_results['documents'])
                    
                    # Safeguard: Limit to ~25,000 characters to avoid exceeding context window
                    if len(full_text) <= 25000:
                        contexts.append({
                            "id": best_reg_id,
                            "text": full_text,
                            "jenis": best_top_meta.get('jenis', ''),
                            "nomor": best_top_meta.get('nomor', ''),
                            "sektor": best_top_meta.get('sektor', ''),
                            "judul": best_top_meta.get('judul', '')
                        })
                        return contexts
                    # If it exceeds 25,000 chars, we DO NOT return contexts here. 
                    # We let it fall through to the standard chunk-based RAG below so the LLM gets the most relevant snippets instead of just the first 25k chars.
                    
    # Fallback to standard chunk-based RAG
    results = collection.query(
        query_texts=[query],
        n_results=overfetch_n,
        where=where_clause
    )
    
    if results and results['documents']:
        for i in range(len(results['documents'][0])):
            meta = results['metadatas'][0][i]
            reg_id = meta.get("reg_id")
            
            # Apply Security Filter (FR-17)
            if not is_allowed(reg_id):
                continue
                
            doc = results['documents'][0][i]
            id_str = results['ids'][0][i]
            
            contexts.append({
                "id": id_str,
                "text": doc,
                "jenis": meta.get('jenis', ''),
                "nomor": meta.get('nomor', ''),
                "sektor": meta.get('sektor', ''),
                "judul": meta.get('judul', '')
            })
            
            # Stop once we have enough authorized chunks
            if len(contexts) >= n_results:
                break
            
    return contexts

def retrieve_graph_contexts(query: str, current_user: dict, max_nodes=10) -> str:
    """FR-16: Queries the SQLite Knowledge Graph based on query keywords and returns a formatted graph string."""
    import sqlite3
    import os
    import re
    
    # 1. Clean query to extract keywords
    stopwords = {"apa", "siapa", "kapan", "dimana", "mengapa", "bagaimana", "dan", "atau", "di", "ke", "dari", "yang", "untuk", "dengan", "tentang", "terkait", "saja", "apakah"}
    words = re.findall(r'\b\w+\b', query.lower())
    keywords = [w for w in words if w not in stopwords and len(w) > 3]
    
    if not keywords:
        return ""
        
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    try:
        conn = sqlite3.connect(db_path)
        c = conn.cursor()
        
        # 2. Find matching nodes
        query_conditions = " OR ".join(["label LIKE ?" for _ in keywords])
        params = [f"%{kw}%" for kw in keywords]
        
        c.execute(f"SELECT id, label, type FROM kg_nodes WHERE {query_conditions} LIMIT {max_nodes}", params)
        matched_nodes = c.fetchall()
        
        if not matched_nodes:
            conn.close()
            return ""
            
        matched_node_ids = [row[0] for row in matched_nodes]
        
        # 3. Find 1-degree connections
        placeholders = ",".join(["?"] * len(matched_node_ids))
        c.execute(f"""
            SELECT e.source_id, n1.label, e.relation, e.target_id, n2.label 
            FROM kg_edges e
            LEFT JOIN kg_nodes n1 ON e.source_id = n1.id
            LEFT JOIN kg_nodes n2 ON e.target_id = n2.id
            WHERE e.source_id IN ({placeholders}) OR e.target_id IN ({placeholders})
            LIMIT 40
        """, matched_node_ids + matched_node_ids)
        
        edges = c.fetchall()
        conn.close()
        
        if not edges:
            return ""
            
        # 4. Format into natural text
        graph_text = "STRUKTUR KNOWLEDGE GRAPH (Entitas dan Relasi yang relevan dengan pertanyaan):\n"
        for src_id, src_label, relation, tgt_id, tgt_label in edges:
            s_lbl = src_label or src_id
            t_lbl = tgt_label or tgt_id
            graph_text += f"- [{s_lbl}] {relation} [{t_lbl}]\n"
            
        return graph_text
    except Exception as e:
        print(f"Graph retrieval error: {e}")
        return ""

@app.post("/api/chat", response_model=ChatResponse)
async def chat_endpoint(req: ChatRequest, current_user: dict = Depends(auth.get_current_user)):
    # FR-24: Admin cannot access document content or search
    if current_user.get("role", "pengguna").lower() == "admin":
        raise HTTPException(status_code=403, detail="Admin sistem tidak memiliki kewenangan untuk mengakses konten dokumen.")
    if not req.messages:
        raise HTTPException(status_code=400, detail="Messages array cannot be empty")
        
    last_user_message = next((m["content"] for m in reversed(req.messages) if m["role"] == "user"), None)
    if not last_user_message:
        raise HTTPException(status_code=400, detail="Missing user message")

    # 1. Retrieve Semantic Context
    contexts = retrieve_contexts(last_user_message, current_user=current_user)
    
    # FR-16: 1.5 Retrieve Graph Context
    graph_context = retrieve_graph_contexts(last_user_message, current_user=current_user)
    
    # Format context for the prompt
    context_str = ""
    sources_to_return = []
    for i, c in enumerate(contexts):
        sources_to_return.append(Source(
            id=c['id'],
            jenis=c['jenis'],
            nomor=c['nomor'],
            sektor=c['sektor'],
            judul=c['judul'],
            snippet=c['text']
        ))
        context_str += f"SUMBER [{i+1}]: {c['jenis']} Nomor {c['nomor']}\nTEKS:\n{c['text']}\n\n"

    if graph_context:
        context_str += f"\n{graph_context}\n\n"

    # Fetch user's contract monitor stats
    user_stats_str = ""
    try:
        conn = auth.get_db_connection()
        c = conn.cursor()
        c.execute("SELECT id FROM compliance_history WHERE user_id = ?", (current_user["id"],))
        rows = c.fetchall()
        total_docs = len(rows)
        conn.close()
        user_stats_str = f"INFO SISTEM (DASHBOARD PENGGUNA): Pengguna saat ini memiliki total {total_docs} dokumen yang tersimpan dan dipantau di dalam Contract Monitor.\n\n"
    except Exception as e:
        print(f"Failed to fetch user stats for chatbot: {e}")

    system_prompt = (
        "Anda adalah pakar hukum teknologi dan penasihat kontrak di Indonesia (Fokus: UU ITE, POJK, KUHPerdata, UU PDP, UU HAM). "
        "Gunakan HANYA konteks hukum yang diberikan untuk menjawab pertanyaan pengguna.\n\n"
        f"{user_stats_str}"
        "PANDUAN KETAT (WAJIB DIIKUTI):\n"
        "1. ZERO META-LANGUAGE: Jangan pernah menggunakan frasa seperti 'Berdasarkan konteks yang diberikan', 'Meskipun konteks tidak menyebutkan', atau 'Menurut konteks hukum'. Anggap fakta hukum sebagai pengetahuan bawaan Anda. Jawab langsung dengan percaya diri layaknya penasihat hukum sungguhan.\n"
        "2. FORMAT PERCAKAPAN & MUDAH DIBACA: Jangan gunakan judul kaku seperti 'Analisis Hukum' atau 'Kesimpulan'. Gunakan alur percakapan yang alami, empatik, dan profesional. Awali jawaban secara langsung dengan 'Ya', 'Tidak', atau 'Tergantung'. Gunakan paragraf pendek, teks tebal (bold) untuk istilah kunci, dan poin-poin agar mudah dibaca di layar ponsel.\n"
        "3. ARGUMEN HUKUM KRITIS:\n"
        "   - Jika ditanya tentang kontrak elektronik vs kertas bermeterai: Tegaskan bahwa meterai hanyalah pajak dokumen, BUKAN syarat sahnya perjanjian. Sahnya perjanjian murni didasarkan pada Pasal 1320 KUHPerdata dan Pasal 5 UU ITE.\n"
        "   - Jika ditanya tentang ancaman penjara untuk utang: Anda WAJIB mengutip 'Pasal 19 ayat (2) UU No. 39 Tahun 1999 tentang Hak Asasi Manusia (UU HAM)' yang melarang keras hukuman pidana/penjara untuk masalah utang piutang perdata. Bedakan dengan jelas antara gagal bayar perdata (wanprestasi) dan niat jahat (penipuan, Pasal 378 KUHP).\n\n"
        "Jawablah dalam Bahasa Indonesia yang profesional dan menenangkan, memberikan solusi yang dapat ditindaklanjuti, dan selalu mengutip dasar hukum/pasal yang relevan secara natural."
    )

    enriched_user_prompt = f"PERTANYAAN PENGGUNA:\n{last_user_message}\n\nKONTEKS HUKUM:\n{context_str}"

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user",   "content": enriched_user_prompt}
    ]
    try:
        answer = call_glm(messages, temperature=0.1, timeout=60)
    except HTTPException as e:
        print(f"Chat fallback triggered due to GLM error: {e.detail}")
        answer = (
            "Maaf, layanan AI sedang mengalami gangguan koneksi sementara. "
            "Silakan coba kirim pertanyaan yang sama beberapa detik lagi."
        )
    # FR-25: Audit log the search query
    log_audit(current_user.get("id", ""), "SEARCH", "", f"Query: {last_user_message[:200]}")
    return ChatResponse(answer=answer, sources=sources_to_return)

@app.post("/api/chat-session")
async def save_chat_session(req: ChatRequest, current_user: dict = Depends(auth.get_current_user)):
    """Save an entire chat session history (called after chat_endpoint or independently)"""
    # For a full implementation, the frontend would pass a session ID, or we generate one.
    # To keep it simple, we will just use a session_id logic here.
    pass # Implementation details added below...


@app.get("/api/pdf/{filename}")
async def serve_pdf(filename: str):
    """
    Return PDF bytes encoded as base64 JSON.
    IDM cannot intercept this because the Content-Type is application/json, not application/pdf.
    The frontend decodes the base64 and creates a blob:// URL to render in an iframe.
    """
    import base64
    safe_filename = os.path.basename(filename)  # prevent path traversal
    file_path = os.path.join(PDFS_DIR, safe_filename)
    if not os.path.isfile(file_path):
        raise HTTPException(status_code=404, detail=f"PDF '{safe_filename}' not found")
    with open(file_path, "rb") as f:
        pdf_bytes = f.read()
    return {"filename": safe_filename, "data": base64.b64encode(pdf_bytes).decode("utf-8")}

@app.get("/api/repository")
async def get_repository(current_user: dict = Depends(auth.get_current_user)):
    # FR-24: Admin cannot access document repository
    if current_user.get("role", "pengguna").lower() == "admin":
        raise HTTPException(status_code=403, detail="Admin sistem tidak memiliki kewenangan untuk mengakses repositori dokumen.")
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    if not os.path.exists(db_path):
        return {"documents": []}
        
    user_id = current_user.get("id")
    role_level = auth.get_role_level(current_user.get("role", "pengguna"))
    
    allowed_klasifikasi = ["Umum"]
    if role_level >= 2:
        allowed_klasifikasi.append("Rahasia")
    if role_level >= 3:
        allowed_klasifikasi.append("Terbatas")
        
    klas_str = ", ".join(f"'{k}'" for k in allowed_klasifikasi)
        
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    
    # Query regulations where klasifikasi is allowed OR explicitly granted
    query = f"""
        SELECT id, judul, nomor, jenis, sektor, status, local_path, klasifikasi 
        FROM regulations 
        WHERE local_path IS NOT NULL AND local_path != ''
        AND status != 'Menunggu Konfirmasi'
        AND (
            klasifikasi IN ({klas_str}) 
            OR id IN (
                SELECT doc_id FROM access_grants 
                WHERE granted_to = ? 
                AND (expires_at IS NULL OR expires_at = '' OR expires_at >= datetime('now'))
            )
        )
    """
    c.execute(query, (user_id,))
    records = c.fetchall()
    
    docs = []
    for row in records:
        # Depending on schema, klasifikasi might be at index 7. Handle safely.
        reg_id, judul, nomor, jenis, sektor, status, local_path = row[:7]
        klasifikasi = row[7] if len(row) > 7 else "Umum"
        filename = os.path.basename(local_path) if local_path else None
        docs.append({
            "id": str(reg_id) if reg_id is not None else None,
            "judul": str(judul) if judul else "",
            "nomor": str(nomor) if nomor is not None else "",
            "jenis": str(jenis) if jenis else "",
            "sektor": str(sektor) if sektor else "",
            "status": str(status) if status else "",
            "klasifikasi": str(klasifikasi) if klasifikasi else "Umum",
            "filename": filename
        })
    conn.close()
    
    return {"documents": docs}


@app.get("/api/repository/pending")
async def get_pending_repository(current_user: dict = Depends(auth.require_role("sekretaris perusahaan"))):
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    c = conn.cursor()
    
    query = """
        SELECT id, judul, nomor, jenis, sektor, status, local_path, klasifikasi 
        FROM regulations 
        WHERE status = 'Menunggu Konfirmasi'
    """
    c.execute(query)
    records = c.fetchall()
    
    docs = []
    for row in records:
        reg_id, judul, nomor, jenis, sektor, status, local_path = row[:7]
        klasifikasi = row[7] if len(row) > 7 else "Umum"
        filename = os.path.basename(local_path) if local_path else None
        docs.append({
            "id": str(reg_id) if reg_id is not None else None,
            "judul": str(judul) if judul else "",
            "nomor": str(nomor) if nomor is not None else "",
            "jenis": str(jenis) if jenis else "",
            "sektor": str(sektor) if sektor else "",
            "status": str(status) if status else "",
            "klasifikasi": str(klasifikasi) if klasifikasi else "Umum",
            "filename": filename
        })
    conn.close()
    
    return {"documents": docs}

@app.post("/api/documents/{doc_id}/grant-access")
async def grant_document_access(
    doc_id: str,
    req: AccessGrantRequest,
    current_user: dict = Depends(auth.get_current_user)
):
    user_role = current_user.get("role", "pengguna").lower()
    user_level = auth.get_role_level(user_role)
    # FR-21 & FR-22: Only Manajer (level 2+) can grant access
    if user_level < 2:
        raise HTTPException(status_code=403, detail="Hanya Manajer, Direktur, atau Sekretaris Perusahaan yang dapat memberikan akses.")
        
    import sqlite3
    import uuid
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    c = conn.cursor()
    
    c.execute("SELECT klasifikasi, judul FROM regulations WHERE id = ?", (doc_id,))
    doc = c.fetchone()
    if not doc:
        conn.close()
        raise HTTPException(status_code=404, detail="Dokumen tidak ditemukan")
        
    klasifikasi = doc[0] if doc[0] else "Umum"
    doc_judul = doc[1] if doc[1] else doc_id
    
    # FR-22: Manajer cannot grant access to Terbatas documents
    if klasifikasi == "Terbatas" and user_level < 3:
        conn.close()
        raise HTTPException(status_code=403, detail="Manajer tidak dapat memberikan akses untuk dokumen Terbatas. Hanya Direktur atau Sekretaris Perusahaan yang berwenang.")
        
    if not req.reason or len(req.reason.strip()) < 5:
        conn.close()
        raise HTTPException(status_code=400, detail="Alasan wajib diisi (minimal 5 karakter).")
        
    grant_id = str(uuid.uuid4())
    c.execute('''INSERT INTO access_grants (id, doc_id, granted_by, granted_to, reason, expires_at)
                 VALUES (?, ?, ?, ?, ?, ?)''',
              (grant_id, doc_id, current_user["id"], req.granted_to, req.reason, req.expires_at))
    conn.commit()
    conn.close()
    
    # FR-25: Audit log the grant
    log_audit(current_user.get("id", ""), "GRANT_ACCESS", doc_id, 
              f"Diberikan kepada: {req.granted_to}, Dokumen: {doc_judul}, Alasan: {req.reason}")
    
    return {"message": "Akses berhasil diberikan."}

@app.get("/api/repository/grants")
async def get_all_grants(current_user: dict = Depends(auth.get_current_user)):
    """FR-23: Sekretaris Perusahaan can see all grants; others see only their own."""
    user_role = current_user.get("role", "pengguna").lower()
    user_level = auth.get_role_level(user_role)
    if user_level < 2:
        raise HTTPException(status_code=403, detail="Akses ditolak.")
    
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    conn.row_factory = sqlite3.Row
    c = conn.cursor()
    
    # FR-23: Sekretaris Perusahaan sees everything; others see only what they granted
    if user_level >= 5:  # Sekretaris Perusahaan
        c.execute('''SELECT ag.id, ag.doc_id, ag.granted_by, ag.granted_to, ag.reason, 
                            ag.expires_at, ag.created_at, r.judul, r.klasifikasi
                     FROM access_grants ag
                     LEFT JOIN regulations r ON ag.doc_id = r.id
                     ORDER BY ag.created_at DESC''')
    else:
        c.execute('''SELECT ag.id, ag.doc_id, ag.granted_by, ag.granted_to, ag.reason, 
                            ag.expires_at, ag.created_at, r.judul, r.klasifikasi
                     FROM access_grants ag
                     LEFT JOIN regulations r ON ag.doc_id = r.id
                     WHERE ag.granted_by = ?
                     ORDER BY ag.created_at DESC''', (current_user["id"],))
    
    rows = c.fetchall()
    conn.close()
    return {"grants": [dict(r) for r in rows]}

@app.delete("/api/repository/grant/{grant_id}")
async def revoke_grant(grant_id: str, current_user: dict = Depends(auth.get_current_user)):
    """FR-23: Sekretaris Perusahaan can revoke any grant."""
    user_level = auth.get_role_level(current_user.get("role", "pengguna"))
    if user_level < 5:  # Only Sekretaris Perusahaan
        raise HTTPException(status_code=403, detail="Hanya Sekretaris Perusahaan yang dapat mencabut pemberian akses.")
    
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    c = conn.cursor()
    c.execute("SELECT doc_id, granted_to FROM access_grants WHERE id = ?", (grant_id,))
    row = c.fetchone()
    if not row:
        conn.close()
        raise HTTPException(status_code=404, detail="Grant tidak ditemukan.")
    
    c.execute("DELETE FROM access_grants WHERE id = ?", (grant_id,))
    conn.commit()
    conn.close()
    
    log_audit(current_user.get("id", ""), "REVOKE_ACCESS", row[0], f"Akses dicabut dari: {row[1]}")
    return {"message": "Akses berhasil dicabut."}

def process_document_background(file_path: str, doc_id: str, filename: str, nomor: str, jenis: str, sektor: str, status: str, klasifikasi: str):
    try:
        parser = LegalDocumentParser()
        full_text = parser.parse_pdf(file_path)
        
        # ── AI Recommendation (FR-4) ──────────────────────────────────────
        messages = [
            {"role": "system", "content": "Anda adalah analis regulasi korporat. Tugas Anda adalah memberikan rekomendasi tingkat kerahasiaan dokumen berdasarkan isinya. Balas hanya dengan satu kata: 'Umum', 'Rahasia', atau 'Terbatas'."},
            {"role": "user", "content": f"Teks dokumen:\n{full_text[:3000]}\n\nBerdasarkan teks ini, rekomendasikan klasifikasi: Umum, Rahasia, atau Terbatas."}
        ]
        
        recommended_klasifikasi = "Umum"
        try:
            raw_content = call_glm(messages, temperature=0.1, timeout=30)
            raw_content = raw_content.lower()
            if "terbatas" in raw_content:
                recommended_klasifikasi = "Terbatas"
            elif "rahasia" in raw_content:
                recommended_klasifikasi = "Rahasia"
        except Exception as llm_err:
            print(f"LLM classification error: {llm_err}")
            
        import sqlite3
        db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
        conn = sqlite3.connect(db_path, timeout=30.0)
        c = conn.cursor()
        c.execute("UPDATE regulations SET status = 'Menunggu Konfirmasi', klasifikasi = ? WHERE id = ?", (recommended_klasifikasi, doc_id))
        conn.commit()
        conn.close()
        print(f"Document {doc_id} set to Pending Confirmation with AI Recommendation: {recommended_klasifikasi}")
        
    except Exception as e:
        print(f"Error in process_document_background: {e}")
        try:
            import sqlite3
            conn = sqlite3.connect(os.path.join(BASE_DIR, "data", "legal_metadata.db"), timeout=30.0)
            c = conn.cursor()
            c.execute("UPDATE regulations SET status = 'Gagal - Error' WHERE id = ?", (doc_id,))
            conn.commit()
            conn.close()
        except:
            pass

def ingest_document_background(file_path: str, doc_id: str, filename: str, nomor: str, jenis: str, sektor: str, status: str, klasifikasi: str):
    try:
        parser = LegalDocumentParser()
        chunker = LegalChunker()
        
        full_text = parser.parse_pdf(file_path)
        
        # ── Duplicate Detection (FR-5) ──────────────────────────────────────
        fingerprint_text = full_text[:1500]
        collection = get_chroma_collection()
        dup_results = collection.query(
            query_texts=[fingerprint_text],
            n_results=1
        )
        
        db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
        conn = sqlite3.connect(db_path, timeout=30.0)
        c = conn.cursor()
        
        if dup_results and dup_results['distances'] and len(dup_results['distances'][0]) > 0:
            dist = dup_results['distances'][0][0]
            if dist < 0.15:
                # Cleanup temp file
                if os.path.exists(file_path):
                    os.remove(file_path)
                c.execute("UPDATE regulations SET status = 'Gagal - Duplikat' WHERE id = ?", (doc_id,))
                conn.commit()
                conn.close()
                return

        print(f"File saved to DB. Now parsing and embedding: {filename}")
        
        # Vector DB Injection
        base_metadata = {
            "reg_id": doc_id,
            "judul": filename.replace('.pdf', ''),
            "nomor": nomor,
            "jenis": jenis,
            "sektor": sektor,
            "status": status
        }
        
        chunks = chunker.chunk_document(full_text, base_metadata)
        
        documents = []
        metadatas = []
        ids = []
        
        import hashlib
        for i, c_data in enumerate(chunks):
            clean_metadata = {k: v for k, v in c_data["metadata"].items() if v is not None}
            chunk_id = f"{nomor}_chunk_{i}"
            hash_id = hashlib.md5(chunk_id.encode('utf-8')).hexdigest()
            
            documents.append(c_data["text"])
            metadatas.append(clean_metadata)
            ids.append(hash_id)
            
        if documents:
            collection = get_chroma_collection()
            collection.add(
                documents=documents,
                metadatas=metadatas,
                ids=ids
            )
            print(f"Successfully added {len(documents)} chunks to ChromaDB!")
        
        c.execute("UPDATE regulations SET status = 'Berlaku' WHERE id = ?", (doc_id,))
        conn.commit()
        conn.close()

        # ── Knowledge Graph Extraction (real-time, FR-KG) ───────────────────
        judul = filename.replace('.pdf', '')
        try:
            extract_and_store_graph(doc_id, full_text, nomor, judul, jenis)
        except Exception as kg_err:
            print(f"[KG] Non-fatal extraction error for {nomor}: {kg_err}")
        return
        
    except Exception as e:
        print(f"Error processing PDF: {e}")
        try:
            conn = sqlite3.connect(os.path.join(BASE_DIR, "data", "legal_metadata.db"), timeout=30.0)
            c = conn.cursor()
            c.execute("UPDATE regulations SET status = 'Gagal - Error' WHERE id = ?", (doc_id,))
            conn.commit()
            conn.close()
        except:
            pass

# ─────────────────────────────────────────────────────────────────────────────
# Knowledge Graph Endpoints
# ─────────────────────────────────────────────────────────────────────────────

@app.get("/api/knowledge-graph")
async def get_knowledge_graph(
    search: str = "",
    node_type: str = "",
    current_user: dict = Depends(auth.require_role("manajer"))
):
    """Returns all KG nodes and edges for the graph visualizer."""
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    conn.row_factory = sqlite3.Row
    c = conn.cursor()

    node_query = "SELECT id, label, type, doc_id FROM kg_nodes WHERE 1=1"
    params = []
    if search:
        node_query += " AND label LIKE ?"
        params.append(f"%{search}%")
    if node_type:
        node_query += " AND type = ?"
        params.append(node_type)
    node_query += " LIMIT 500"

    c.execute(node_query, params)
    nodes = [dict(r) for r in c.fetchall()]

    node_ids = {n["id"] for n in nodes}

    # Only return edges where both endpoints are in the node set
    c.execute("SELECT id, source_id, target_id, relation, doc_id FROM kg_edges LIMIT 2000")
    all_edges = c.fetchall()
    edges = [dict(e) for e in all_edges if e["source_id"] in node_ids and e["target_id"] in node_ids]

    c.execute("SELECT COUNT(*) FROM kg_nodes")
    total_nodes = c.fetchone()[0]
    c.execute("SELECT COUNT(*) FROM kg_edges")
    total_edges = c.fetchone()[0]

    conn.close()
    return {"nodes": nodes, "edges": edges, "total_nodes": total_nodes, "total_edges": total_edges}


def _rebuild_kg_batch():
    """Background worker: rebuilds KG by pulling text from ChromaDB chunks.
    Used for documents ingested via the scraper that have no local PDF files.
    """
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    conn.row_factory = sqlite3.Row
    c = conn.cursor()
    c.execute("SELECT id, nomor, jenis, judul FROM regulations WHERE status LIKE 'Berlaku%' AND id NOT IN (SELECT DISTINCT doc_id FROM kg_nodes WHERE doc_id IS NOT NULL)")
    docs = c.fetchall()
    conn.close()

    collection = get_chroma_collection()
    print(f"[KG Rebuild] Starting ChromaDB-based batch for {len(docs)} documents...")
    success, failed, skipped = 0, 0, 0

    for doc in docs:
        try:
            # Query ChromaDB for chunks that belong to this regulation
            results = collection.query(
                query_texts=[doc["nomor"] or doc["judul"] or ""],
                n_results=5,
                where={"reg_id": str(doc["id"])}
            )
            documents_list = results.get("documents", [[]])[0]
            if not documents_list:
                # fallback: try matching by nomor in metadata
                skipped += 1
                continue

            full_text = "\n\n".join(documents_list)
            extract_and_store_graph(
                doc["id"], full_text,
                doc["nomor"] or "", doc["judul"] or "", doc["jenis"] or ""
            )
            success += 1
        except Exception as e:
            print(f"[KG Rebuild] Failed for {doc['nomor']}: {e}")
            failed += 1

    print(f"[KG Rebuild] Done. Success: {success}, Skipped (no chunks): {skipped}, Failed: {failed}")


class ScenarioAnalyzeRequest(BaseModel):
    scenario: str

@app.post("/api/knowledge-graph/analyze-scenario")
async def analyze_kg_scenario(req: ScenarioAnalyzeRequest):
    """
    Uses LLM to dynamically select which node IDs are relevant to the requested scenario.
    """
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    c = conn.cursor()
    
    # Fetch all node IDs and labels
    c.execute("SELECT id, label FROM kg_nodes")
    nodes = c.fetchall()
    conn.close()
    
    # We only send a subset of data to avoid exceeding context if it's too large,
    # but 1400 nodes is about ~50k chars which is perfectly fine for modern LLMs.
    nodes_str = "\n".join([f"ID: {n[0]} | Label: {n[1]}" for n in nodes])
    
    messages = [
        {"role": "system", "content": "You are an expert legal knowledge graph analyst. Your job is to return a JSON array of Node IDs that are highly relevant to the user's requested scenario. Be strict and only return nodes directly involved with the scenario."},
        {"role": "user", "content": f"Here is the list of all nodes in our knowledge graph:\n\n{nodes_str}\n\nScenario: {req.scenario}\n\nReturn ONLY a JSON array of strings containing the exact IDs of the nodes that are highly relevant to this scenario. Example: [\"node1\", \"node2\"]. Return nothing else."}
    ]
    
    try:
        raw_response = call_glm(messages, temperature=0.1, timeout=90)
        
        # Parse out JSON block
        import re
        import json
        match = re.search(r'\[.*?\]', raw_response, re.DOTALL)
        if match:
            node_ids = json.loads(match.group(0))
            return {"status": "success", "matchedNodeIds": node_ids}
        else:
            return {"status": "error", "matchedNodeIds": []}
    except Exception as e:
        print(f"LLM Scenario Error: {e}")
        return {"status": "error", "matchedNodeIds": []}

@app.post("/api/knowledge-graph/rebuild")
async def rebuild_knowledge_graph(
    background_tasks: BackgroundTasks,
    current_user: dict = Depends(auth.require_role("admin"))
):
    """Admin-only: triggers batch KG rebuild for all existing documents."""
    background_tasks.add_task(_rebuild_kg_batch)
    return {"message": "Rebuild dimulai di background. Proses ini bisa memakan waktu 30-60 menit."}


@app.delete("/api/knowledge-graph/document/{doc_id}")
async def delete_doc_from_graph(
    doc_id: str,
    current_user: dict = Depends(auth.require_role("manajer"))
):
    """Removes all KG nodes and edges created by a specific document."""
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    c = conn.cursor()
    c.execute("DELETE FROM kg_edges WHERE doc_id = ?", (doc_id,))
    c.execute("DELETE FROM kg_nodes WHERE doc_id = ?", (doc_id,))
    conn.commit()
    conn.close()
    return {"message": "Data graph untuk dokumen ini telah dihapus."}

@app.get("/api/admin/dashboard")
async def admin_dashboard(current_user: dict = Depends(auth.require_role("admin"))):
    """FR-26: Admin-only dashboard with system stats and audit logs."""
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    conn.row_factory = sqlite3.Row
    c = conn.cursor()

    # 1. Document Processing Status & Details
    c.execute("SELECT id, judul, status, nomor FROM regulations")
    all_docs = c.fetchall()
    
    doc_details = {
        "Berlaku": [],
        "Tidak Berlaku": [],
        "Memproses": [],
        "Gagal": []
    }
    
    for d in all_docs:
        s = (d["status"] or "").strip()
        doc_obj = {"id": d["id"], "judul": d["judul"], "nomor": d["nomor"], "status": s}
        
        if s == "Memproses":
            doc_details["Memproses"].append(doc_obj)
        elif s.startswith("Gagal"):
            doc_details["Gagal"].append(doc_obj)
        elif s == "Tidak Berlaku" or "Dicabut" in s:
            doc_details["Tidak Berlaku"].append(doc_obj)
        else:
            doc_details["Berlaku"].append(doc_obj)
            
    doc_status = {
        "Berlaku": len(doc_details["Berlaku"]),
        "Tidak Berlaku": len(doc_details["Tidak Berlaku"]),
        "Memproses": len(doc_details["Memproses"]),
        "Gagal": len(doc_details["Gagal"])
    }

    # 2. Document Volume by Klasifikasi
    c.execute("SELECT klasifikasi, COUNT(*) as count FROM regulations WHERE klasifikasi IS NOT NULL GROUP BY klasifikasi")
    klas_rows = c.fetchall()
    doc_by_klasifikasi = {r["klasifikasi"]: r["count"] for r in klas_rows}

    # 3. Document Volume by Jenis
    c.execute("SELECT jenis, COUNT(*) as count FROM regulations GROUP BY jenis ORDER BY count DESC LIMIT 10")
    jenis_rows = c.fetchall()
    doc_by_jenis = [dict(r) for r in jenis_rows]

    # 4. Active Access Grants count
    c.execute("SELECT COUNT(*) FROM access_grants WHERE expires_at IS NULL OR expires_at = '' OR expires_at >= datetime('now')")
    active_grants = c.fetchone()[0]

    # 5. Recent Audit Logs (last 100)
    c.execute("SELECT id, timestamp, user_id, action_type, resource_id, details FROM audit_logs ORDER BY id DESC LIMIT 100")
    audit_rows = c.fetchall()
    audit_logs = [dict(r) for r in audit_rows]

    conn.close()

    # 6. System Health
    chroma_ok = False
    try:
        col = get_chroma_collection()
        chroma_ok = col.count() >= 0
    except:
        pass

    return {
        "doc_status": doc_status,
        "doc_details": doc_details,
        "doc_by_klasifikasi": doc_by_klasifikasi,
        "doc_by_jenis": doc_by_jenis,
        "active_grants": active_grants,
        "audit_logs": audit_logs,
        "system_health": {
            "sqlite": True,
            "chromadb": chroma_ok
        }
    }

@app.delete("/api/repository/document/{doc_id}")
async def delete_document(doc_id: str, current_user: dict = Depends(auth.require_role("sekretaris perusahaan"))):
    """Deletes a document from the repository."""
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    
    # Check if document exists
    c.execute("SELECT id, local_path FROM regulations WHERE id = ?", (doc_id,))
    doc = c.fetchone()
    if not doc:
        conn.close()
        raise HTTPException(status_code=404, detail="Dokumen tidak ditemukan")
    
    local_path = doc[1]
    
    # Delete physical file
    if local_path and os.path.exists(local_path):
        try:
            os.remove(local_path)
        except Exception as e:
            print(f"Failed to delete file {local_path}: {e}")
            
    # Delete from DB
    c.execute("DELETE FROM regulations WHERE id = ?", (doc_id,))
    c.execute("DELETE FROM access_grants WHERE doc_id = ?", (doc_id,))
    c.execute("DELETE FROM kg_edges WHERE doc_id = ?", (doc_id,))
    c.execute("DELETE FROM kg_nodes WHERE doc_id = ?", (doc_id,))
    
    conn.commit()
    conn.close()
    
    # Delete from ChromaDB
    try:
        collection = get_chroma_collection()
        collection.delete(where={"reg_id": doc_id})
    except Exception as e:
        print(f"Failed to delete from ChromaDB: {e}")
        
    return {"message": "Dokumen berhasil dihapus"}

@app.delete("/api/repository/failed")
async def delete_failed_documents(current_user: dict = Depends(auth.require_role("manajer"))):
    """Deletes all documents that failed processing (e.g., Duplicates)."""
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path, timeout=30.0)
    c = conn.cursor()
    
    c.execute("SELECT id, local_path FROM regulations WHERE status LIKE 'Gagal%'")
    failed_docs = c.fetchall()
    
    deleted_count = 0
    for doc in failed_docs:
        doc_id, local_path = doc
        # Delete physical file
        if local_path and os.path.exists(local_path):
            try:
                os.remove(local_path)
            except Exception as e:
                print(f"Error deleting file {local_path}: {e}")
        # Delete from DB
        c.execute("DELETE FROM regulations WHERE id = ?", (doc_id,))
        deleted_count += 1
        
    conn.commit()
    conn.close()
    
    return {"message": f"Berhasil menghapus {deleted_count} dokumen yang gagal."}

@app.post("/api/upload")
async def upload_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...), 
    doc_type: str = Form("regulations"),
    klasifikasi: str = Form("Umum"),
    current_user: dict = Depends(auth.require_role("manajer"))
):
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
        
    import sqlite3
    import uuid
    from datetime import datetime
    
    # Setup directories
    pdfs_dir = os.path.join(BASE_DIR, "data", "pdfs")
    os.makedirs(pdfs_dir, exist_ok=True)
    
    file_path = os.path.join(pdfs_dir, file.filename)
    
    # Save physical file temporarily
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        # ── Proceed with Saving Metadata ─────────────────────────────────────
        db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
        conn = sqlite3.connect(db_path, timeout=30.0)
        c = conn.cursor()
        c.execute('''CREATE TABLE IF NOT EXISTS regulations
                     (id INTEGER PRIMARY KEY AUTOINCREMENT, domain TEXT, judul TEXT, 
                      nomor TEXT, jenis TEXT, sektor TEXT, status TEXT, 
                      detail_url TEXT, download_url TEXT, local_path TEXT, klasifikasi TEXT)''')
                      
        # Generate ID and Metadata
        doc_id = str(uuid.uuid4())[:8]
        judul = file.filename.replace('.pdf', '')
        nomor = f"CUSTOM-{doc_id}"
        
        if jenis_dokumen:
            jenis = jenis_dokumen
        elif doc_type == "internal":
            jenis = "Dokumen Internal"
            sektor = "Dokumen Internal"
        else:
            jenis = "Regulasi Custom"
            sektor = "Upload Manual"
        
        status = "Memproses" # Initialize with processing status
        
        c.execute('''INSERT OR REPLACE INTO regulations 
                     (judul, download_url, nomor, jenis, sektor, status, detail_url, local_path, klasifikasi, domain) 
                     VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)''',
                  (judul, "", nomor, jenis, sektor, status, "", file_path, klasifikasi, "Custom"))
        # Get the actual auto-incremented ID
        doc_id = str(c.lastrowid)
        conn.commit()
        conn.close()
        
        # Enqueue background processing task
        background_tasks.add_task(process_document_background, file_path, doc_id, file.filename, nomor, jenis, sektor, status, klasifikasi)
        
        return {"status": "success", "message": "Dokumen berhasil diunggah dan sedang diproses di latar belakang."}
        
    except Exception as e:
        print(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail="Terjadi kesalahan saat memproses dokumen.")

@app.post("/api/analyze")
async def analyze_document(req: AnalyzeRequest):
    """Analyzes a document from ChromaDB chunks + LLM for overview and status."""
    collection = get_chroma_collection()
    
    # Step 1: Fetch all chunks for this specific document
    # Prefer reg_id (unique) but fall back to nomor if not available
    try:
        if req.reg_id:
            results = collection.get(
                where={"reg_id": req.reg_id},
                include=["documents"]
            )
        else:
            results = collection.get(
                where={"nomor": req.nomor},
                include=["documents"]
            )
        chunks = results.get("documents", [])
    except Exception as e:
        print(f"ChromaDB fetch error: {e}")
        chunks = []
    
    if not chunks:
        return {
            "total_pasal": 0,
            "overview": "Dokumen ini belum memiliki data di knowledge base.",
            "status": {"dicabut": [], "diubah_dengan": []},
            "outline": []
        }
    
    # Step 2: Build outline + find max Pasal number (more accurate than chunk count)
    import re
    outline = []
    seen = set()
    max_pasal_num = 0
    
    for chunk in chunks:
        lines = chunk.strip().split('\n')
        first_line = lines[0].strip() if lines else ""
        
        # Detect BAB headers
        if re.match(r'^BAB\s+[IVXLC\d]+', first_line, re.IGNORECASE):
            entry = first_line[:80]
            if entry not in seen:
                outline.append({"type": "bab", "text": entry})
                seen.add(entry)
        
        # Detect Pasal headers — track the max number found
        pasal_match = re.search(r'Pasal\s+(\d+)', chunk, re.IGNORECASE)
        if pasal_match:
            num = int(pasal_match.group(1))
            if num > max_pasal_num:
                max_pasal_num = num
            
            entry = first_line[:60]
            if re.match(r'^Pasal\s+\d+', first_line, re.IGNORECASE) and entry not in seen:
                outline.append({"type": "pasal", "text": entry})
                seen.add(entry)

    # Step 3: Take sample text for LLM (first 6 chunks ≈ 3500 chars)
    sample_text = "\n---\n".join(chunks[:6])[:3500]
    
    system_prompt = (
        "Anda adalah analis hukum OJK Indonesia. Berikan analisis singkat dan akurat dalam format JSON yang ketat. "
        "Jangan menambahkan teks di luar JSON."
    )
    
    user_prompt = (
        f"Analisis dokumen hukum berikut:\n"
        f"Judul: {req.judul}\nNomor: {req.nomor}\n\n"
        f"TEKS DOKUMEN:\n{sample_text}\n\n"
        f"Kembalikan HANYA JSON berikut (tanpa markdown, tanpa kode blok):\n"
        '{"overview": "ringkasan 2-3 kalimat", '
        '"dicabut": ["nama peraturan yang dicabut jika ada, atau array kosong"], '
        '"diubah_dengan": ["nama peraturan pengubah jika ada, atau array kosong"]}'
    )
    
    messages_llm = [
        {"role": "system", "content": system_prompt},
        {"role": "user",   "content": user_prompt}
    ]
    
    overview = "Tidak dapat mengambil ringkasan dari dokumen ini."
    dicabut = []
    diubah_dengan = []
    
    try:
        raw_content = call_glm(messages_llm, temperature=0.1, timeout=60)
        raw_content = re.sub(r'```json|```', '', raw_content).strip()
        parsed = json.loads(raw_content)
        overview = parsed.get("overview", overview)
        dicabut = parsed.get("dicabut", [])
        diubah_dengan = parsed.get("diubah_dengan", [])
    except HTTPException:
        raise
    except Exception as e:
        print(f"LLM analysis error: {e}")
    
    return {
        "total_pasal": max_pasal_num,
        "overview": overview,
        "status": {"dicabut": dicabut, "diubah_dengan": diubah_dengan},
        "outline": outline
    }

@app.post("/api/analyze-pasals")
async def analyze_pasals(req: AnalyzeRequest):
    """
    Deep analysis: uses LLM to analyze each Pasal individually,
    producing status, perbandingan (comparison), and hubungan (hierarchy).
    """
    collection = get_chroma_collection()

    try:
        if req.reg_id:
            results = collection.get(where={"reg_id": req.reg_id}, include=["documents"])
        else:
            results = collection.get(where={"nomor": req.nomor}, include=["documents"])
        chunks = results.get("documents", [])
    except Exception as e:
        print(f"ChromaDB error: {e}")
        chunks = []

    if not chunks:
        return {"pasals": [], "error": "Tidak ada data untuk dokumen ini."}

    # ── Build pasal map: {pasal_num (int): text (str)} ──
    pasal_texts: dict[int, str] = {}
    for chunk in chunks:
        m = re.search(r'Pasal\s+(\d+)', chunk, re.IGNORECASE)
        if m:
            num = int(m.group(1))
            if num not in pasal_texts:
                pasal_texts[num] = chunk[:500]

    sorted_pasals = sorted(pasal_texts.items())[:15]   # limit to 15 most relevant pasals

    if not sorted_pasals:
        return {"pasals": [], "error": "Tidak dapat mengekstrak Pasal dari dokumen ini."}

    # ── Build LLM prompt ──
    doc_text = ""
    for num, text in sorted_pasals:
        doc_text += f"\nPasal {num}:\n{text[:380]}\n---\n"

    system_prompt = (
        "Anda adalah analis regulasi di Indonesia yang sangat teliti dan akurat. "
        "Selalu jawab dalam Bahasa Indonesia. Jangan menambahkan teks di luar JSON."
    )
    user_prompt = (
        f"Analisis regulasi berikut secara mendalam:\n"
        f"Judul: {req.judul}\nNomor: {req.nomor}\n\n"
        f"TEKS REGULASI:\n{doc_text}\n\n"
        f"Untuk SETIAP Pasal di atas, berikan analisis dalam format JSON strict berikut "
        f"(kembalikan HANYA array JSON, tanpa teks tambahan):\n"
        '[{"pasal":"Pasal 1","status":"Tidak Berubah",'
        '"isi":"ringkasan isi pasal dalam 1-2 kalimat",'
        '"perbandingan":"perbandingan konten pasal ini dengan regulasi sebelumnya yang dicabut atau diubah",'
        '"hubungan":"hubungan dengan Undang-Undang atau Peraturan Pemerintah yang lebih tinggi"},'
        '{"pasal":"Pasal 2",...}]'
        '\n\nNilai status yang valid: "Tidak Berubah", "Diubah", "Baru", "Dicabut".'
    )

    try:
        raw = call_glm([
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_prompt},
        ], temperature=0.1, timeout=180)
        raw = re.sub(r'```json|```', '', raw).strip()

        # The LLM might return an object {"pasals":[...]} OR a bare array [...]
        parsed = json.loads(raw)
        if isinstance(parsed, dict):
            items = parsed.get("pasals", parsed.get("data", []))
        elif isinstance(parsed, list):
            items = parsed
        else:
            items = []

        # Attach the actual extracted pasal text for the frontend to display
        pasal_map = {num: text for num, text in sorted_pasals}
        for item in items:
            m = re.search(r'\d+', item.get("pasal", ""))
            if m:
                item["isi_teks"] = pasal_map.get(int(m.group()), "")[:350]

        return {"pasals": items}

    except json.JSONDecodeError as e:
        print(f"JSON parse error in analyze-pasals: {e}\nRaw: {raw[:300]}")
        return {"pasals": [], "error": "LLM memberikan respons yang tidak valid. Coba lagi."}
    except Exception as e:
        print(f"Deep analysis error: {e}")
        return {"pasals": [], "error": str(e)}


def process_single_pasal(c, doc_type, pihak_1, pihak_2, sektor, pokok, collection, COSINE_THRESHOLD):
    pasal_label = c.get("pasal", "Pasal")
    desc = c.get("deskripsi", "")
    if not desc or len(desc) < 10:
        return None

    # Map doc_type to doc_category
    doc_cat = "NDA" if "NDA" in str(doc_type).upper() or "NON-DISCLOSURE" in str(doc_type).upper() else "PKS"

    # Enrich the search query with document context for better RAG retrieval
    search_query = f"[{doc_type}] [{sektor}] {desc[:400]}"
    sr = collection.query(
        query_texts=[search_query], 
        n_results=3,
        where={"doc_category": {"$in": [doc_cat, "UMUM"]}}
    )

    supporting_regulations = []
    relevant_refs = []

    if sr and sr.get('documents') and len(sr['documents'][0]) > 0:
        for i in range(len(sr['documents'][0])):
            dist = sr['distances'][0][i]
            if dist < COSINE_THRESHOLD:
                doc_text = sr['documents'][0][i]
                meta = sr['metadatas'][0][i]
                reg_name = f"{meta.get('jenis', 'Aturan')} Nomor {meta.get('nomor', 'N/A')}"

                # Relevance confirmation micro-call
                if is_regulation_relevant(desc, doc_text, reg_name):
                    supporting_regulations.append({
                        "jenis": meta.get('jenis', 'Aturan'),
                        "nomor": meta.get('nomor', 'N/A'),
                        "sektor": meta.get('sektor', 'Umum'),
                        "teks": doc_text
                    })
                    relevant_refs.append(f"--- {reg_name} ---\n{doc_text[:700]}")
                else:
                    print(f"  [Relevance] {reg_name} rejected for {pasal_label}")

    if not relevant_refs:
        return {
            "pasal": pasal_label,
            "isi_pasal": desc[:500],
            "status": "TIDAK DIATUR",
            "penjelasan": "Tidak ditemukan regulasi yang secara langsung relevan dengan klausul ini dalam database saat ini.",
            "rekomendasi": "Pastikan klausul ini wajar secara keperdataan dan tidak bertentangan dengan KUH Perdata.",
            "ai_analysis": "",
            "supporting_regulations": supporting_regulations
        }

    combined_regs = "\n\n".join(relevant_refs)
    prompt_verify = (
        "Anda adalah auditor kepatuhan hukum senior Indonesia yang sangat teliti. "
        "Analisis klausul kontrak berikut secara mendalam menggunakan regulasi yang tersedia.\n\n"

        f"KONTEKS DOKUMEN:\n"
        f"- Jenis: {doc_type}\n"
        f"- Pihak Pertama: {pihak_1}\n"
        f"- Pihak Kedua: {pihak_2}\n"
        f"- Sektor: {sektor}\n"
        f"- Pokok Perjanjian: {pokok}\n\n"

        f"KLAUSUL YANG DIANALISIS ({pasal_label}):\n{desc}\n\n"

        f"REGULASI YANG BERLAKU:\n{combined_regs}\n\n"

        "INSTRUKSI ANALISIS (ikuti urutan ini):\n"
        "1. Identifikasi: Apa topik hukum utama klausul ini? (pembayaran, kerahasiaan, PHK, dll)\n"
        "2. Kesesuaian: Apakah klausul ini SECARA EKSPLISIT diatur, dilarang, atau diwajibkan oleh regulasi di atas?\n"
        "   Kutip pasal/ayat spesifik dari regulasi jika ada.\n"
        "3. Risiko: Identifikasi risiko hukum konkret bagi pihak yang lebih lemah.\n"
        "4. Verdict: Berikan satu dari tiga status:\n"
        "   - SESUAI: klausul patuh dan tidak berisiko\n"
        "   - BERESIKO: klausul perlu perbaikan atau klarifikasi\n"
        "   - FATAL: klausul melanggar ketentuan wajib regulasi\n\n"
        "Kembalikan HANYA JSON berikut (tanpa markdown, tanpa teks lain):\n"
        '{\n'
        '  "status": "SESUAI" | "BERESIKO" | "FATAL",\n'
        '  "pasal_regulasi_terkait": "misal: Pasal 15 ayat (3) PP 45/2015",\n'
        '  "penjelasan": "penjelasan spesifik 1-2 kalimat merujuk regulasi",\n'
        '  "rekomendasi": "saran perbaikan konkret (kosongkan jika SESUAI)",\n'
        '  "ai_analysis": "dampak bisnis jangka panjang (isi hanya jika BERESIKO atau FATAL)"\n'
        '}'
    )

    try:
        raw_v = call_glm(
            [{"role": "user", "content": prompt_verify}],
            temperature=0.1,
            timeout=90
        )
        raw_v = re.sub(r'```json|```', '', raw_v).strip()
        ans = json.loads(raw_v)

        status_raw = str(ans.get("status", "BERESIKO")).upper()
        if "SESUAI" in status_raw:       status_final = "SESUAI"
        elif "FATAL" in status_raw:       status_final = "FATAL"
        elif "TIDAK DIATUR" in status_raw: status_final = "TIDAK DIATUR"
        else:                              status_final = "BERESIKO"

        return {
            "pasal": pasal_label,
            "isi_pasal": desc[:500],
            "status": status_final,
            "pasal_regulasi_terkait": ans.get("pasal_regulasi_terkait", ""),
            "penjelasan": ans.get("penjelasan", "Gagal diverifikasi"),
            "rekomendasi": ans.get("rekomendasi", ""),
            "ai_analysis": ans.get("ai_analysis", ""),
            "supporting_regulations": supporting_regulations
        }

    except HTTPException as e:
        print(f"[Compliance] Verification HTTP error: {e.detail}")
        return {
            "pasal": pasal_label, "isi_pasal": desc[:500], "status": "ERROR",
            "pasal_regulasi_terkait": "",
            "penjelasan": "Layanan AI sedang mengalami gangguan koneksi sementara. Silakan coba lagi dalam beberapa saat.",
            "rekomendasi": "Coba ulang analisis.",
            "ai_analysis": "", "supporting_regulations": supporting_regulations
        }
    except Exception as e:
        print(f"[Compliance] Parsing error for {pasal_label}: {e}")
        return {
            "pasal": pasal_label, "isi_pasal": desc[:500], "status": "ERROR",
            "pasal_regulasi_terkait": "",
            "penjelasan": "Format respons LLM tidak valid.",
            "rekomendasi": "-",
            "ai_analysis": "", "supporting_regulations": supporting_regulations
        }

def extract_text_multi_format(file_path: str, filename: str, use_ocr: bool = False) -> str:
    """
    Extracts text from various file formats: PDF, DOCX, XLSX, PPTX, JPG, PNG, TXT.
    Routes to the appropriate extractor based on the file extension.
    """
    ext = filename.lower().split('.')[-1]
    
    if ext == 'pdf':
        if use_ocr:
            print(f"  [OCR] Forcing traditional OCR extraction for {filename}")
            import sys
            if BASE_DIR not in sys.path:
                sys.path.append(BASE_DIR)
            from parser.pdf_parser import LegalDocumentParser
            parser = LegalDocumentParser()
            return parser.parse_pdf(file_path, force_ocr=True)
        else:
            return extract_text_hybrid(file_path, force_vlm=False)
        
    elif ext == 'txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
            
    elif ext in ['jpg', 'jpeg', 'png']:
        if use_ocr:
            print(f"  [OCR] Traditional image extraction for {filename}")
            import sys
            if BASE_DIR not in sys.path:
                sys.path.append(BASE_DIR)
            from parser.pdf_parser import LegalDocumentParser
            parser = LegalDocumentParser()
            if parser.ocr:
                import numpy as np
                from PIL import Image
                img = Image.open(file_path).convert("RGB")
                img_array = np.array(img)
                result = parser.ocr.ocr(img_array, cls=False)
                page_text = []
                if result and result[0]:
                    for line in result[0]:
                        page_text.append(line[1][0])
                return " ".join(page_text)
            else:
                return "[OCR GAGAL INISIALISASI]"
        else:
            print(f"  [VLM] Image extraction for {filename}")
            with open(file_path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode('utf-8')
            vlm_text = call_glm_vision(img_b64, VLM_PAGE_PROMPT, timeout=60)
            return vlm_text.strip()
        
    elif ext == 'docx':
        try:
            import docx
            doc = docx.Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:
                        full_text.append(" | ".join(row_data))
            return "\n".join(full_text)
        except Exception as e:
            print(f"Error parsing DOCX: {e}")
            return ""
            
    elif ext == 'pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        full_text.append(shape.text.strip())
            return "\n".join(full_text)
        except Exception as e:
            print(f"Error parsing PPTX: {e}")
            return ""
            
    elif ext == 'xlsx':
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            ws = wb.active
            full_text = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if row_data:
                    full_text.append(" | ".join(row_data))
            return "\n".join(full_text)
        except Exception as e:
            print(f"Error parsing XLSX: {e}")
            return ""
            
    else:
        return ""


@app.post("/api/check-compliance")
async def check_compliance(file: UploadFile = File(...), use_ocr: str = Form("false")):
    """
    Accepts a document upload (PDF, DOCX, XLSX, PPTX, JPG, PNG, TXT), extracts text, 
    and uses a multi-step LLM pipeline to cross-check clauses against the OJK repository.
    """
    allowed_exts = ['.pdf', '.docx', '.xlsx', '.pptx', '.jpg', '.jpeg', '.png', '.txt']
    ext = os.path.splitext(file.filename.lower())[1]
    
    if ext not in allowed_exts:
        raise HTTPException(status_code=400, detail="Format file tidak didukung. Harap unggah PDF, DOCX, XLSX, PPTX, JPG, PNG, atau TXT.")
        
    upload_id = str(uuid.uuid4())[:8]
    temp_path = os.path.join(PDFS_DIR, f"temp_check_{upload_id}_{file.filename}")
    
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        # ── Pass 0: Multi-Format Extraction ──────────────────────────────────
        print(f"[Compliance] Starting extraction for: {file.filename}")
        is_ocr = use_ocr.lower() == "true"
        full_text = extract_text_multi_format(temp_path, file.filename, use_ocr=is_ocr)
        text_snippet = full_text[:35000]

        if len(text_snippet.strip()) < 100:
            raise HTTPException(
                status_code=400,
                detail="Dokumen tidak dapat dibaca. Coba unggah ulang atau pastikan PDF tidak terenkripsi."
            )

        # ── Pass 0A: Document Metadata Understanding ────────────────────────
        print("[Compliance] Pass 0A: Extracting document metadata...")
        doc_ctx = understand_document(full_text)
        doc_type = doc_ctx.get("jenis_dokumen", "Kontrak")
        pihak_1  = doc_ctx.get("pihak_pertama", "Pihak Pertama")
        pihak_2  = doc_ctx.get("pihak_kedua", "Pihak Kedua")
        sektor   = doc_ctx.get("sektor_bisnis", "umum")
        pokok    = doc_ctx.get("pokok_perjanjian", "")

        # ── Pass 0B: Dedicated signing date extraction ──────────────────────
        print("[Compliance] Pass 0B: Extracting signing date...")
        tanggal_pembuatan = extract_signing_date(full_text)

        # ── Pass 0C: Dedicated duration extraction ──────────────────────────
        print("[Compliance] Pass 0C: Extracting contract duration...")
        durasi_bulan = extract_duration_months(full_text)

        # ── Compute expiry date ─────────────────────────────────────────────
        tanggal_berakhir = None
        if tanggal_pembuatan and durasi_bulan:
            try:
                start_date = datetime.strptime(tanggal_pembuatan, "%Y-%m-%d")
                end_date = start_date + relativedelta(months=int(durasi_bulan))
                tanggal_berakhir = end_date.strftime("%Y-%m-%d")
                print(f"[Compliance] Expiry computed: {tanggal_pembuatan} + {durasi_bulan}mo = {tanggal_berakhir}")
            except Exception as e:
                print(f"[Compliance] Failed to compute expiry: {e}")
                
        # ── Pass 0D: Explicit end date fallback ─────────────────────────────
        if not tanggal_berakhir:
            print("[Compliance] Computation failed/missing, attempting explicit end date extraction...")
            tanggal_berakhir = extract_explicit_end_date(full_text)

        # ── Pass 1: Clause Extraction ───────────────────────────────────────
        print("[Compliance] Pass 1: Extracting clauses...")
        pasal_items = extract_pasal_items(text_snippet, max_items=20, max_chars_per_pasal=1200)

        # Fallback to LLM extractor for non-standard contracts
        if len(pasal_items) < 2:
            print("[Compliance] Regex found <2 clauses, falling back to LLM extractor...")
            pasal_items = llm_extract_pasal_items(full_text, max_items=20)

        # Final fallback: treat whole document as one unit
        if not pasal_items:
            pasal_items = [{"pasal": "Dokumen Keseluruhan", "deskripsi": text_snippet[:1200]}]

        print(f"[Compliance] Found {len(pasal_items)} clauses. Running analysis...")

        # ── Pass 2: Per-clause RAG + LLM Compliance Check ──────────────────
        collection = get_chroma_collection()
        results = []
        COSINE_THRESHOLD = 0.45  # tightened from 0.60 to reduce wrong-domain citations

        def worker(c):
            return process_single_pasal(c, doc_type, pihak_1, pihak_2, sektor, pokok, collection, COSINE_THRESHOLD)

        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            for res in executor.map(worker, pasal_items[:20]):
                if res:
                    results.append(res)

        # ── Pass 4: Document-Level Summary ─────────────────────────────────
        fatal_count    = sum(1 for r in results if r["status"] == "FATAL")
        beresiko_count = sum(1 for r in results if r["status"] == "BERESIKO")
        sesuai_count   = sum(1 for r in results if r["status"] == "SESUAI")
        tidak_diatur_count = sum(1 for r in results if r["status"] == "TIDAK DIATUR")

        # Calculate score only against regulated clauses
        regulated_count = len(results) - tidak_diatur_count
        skor = round((sesuai_count / max(regulated_count, 1)) * 100) if regulated_count > 0 else 100

        summary = {
            "jenis_dokumen": doc_type,
            "pihak_pertama": pihak_1,
            "pihak_kedua": pihak_2,
            "sektor_bisnis": sektor,
            "pokok_perjanjian": pokok,
            "tanggal_berakhir": tanggal_berakhir,
            "total_klausul": len(results),
            "sesuai": sesuai_count,
            "beresiko": beresiko_count,
            "fatal": fatal_count,
            "skor_kepatuhan": skor
        }

        print(f"[Compliance] Done. Score: {summary['skor_kepatuhan']}% | "
              f"SESUAI={sesuai_count} BERESIKO={beresiko_count} FATAL={fatal_count}")

        return {"report": results, "summary": summary}

    except HTTPException:
        raise
    except Exception as e:
        print(f"Compliance checker unexpected error: {e}")
        raise HTTPException(status_code=500, detail="Terjadi kesalahan internal saat memproses dokumen.")
    finally:
        # Automagically clean up uploaded PKS memory footprint
        if os.path.exists(temp_path):
            try: os.remove(temp_path)
            except: pass

class KGExclusionCreate(BaseModel):
    entity_name: str

@app.get("/api/admin/kg-exclusions")
async def get_kg_exclusions(current_user: dict = Depends(auth.require_role("admin"))):
    """FR-30: Get all entity exclusions"""
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute("SELECT id, entity_name, created_at FROM kg_exclusions ORDER BY created_at DESC")
    exclusions = [{"id": row[0], "entity_name": row[1], "created_at": row[2]} for row in c.fetchall()]
    conn.close()
    return {"exclusions": exclusions}

@app.post("/api/admin/kg-exclusions")
async def add_kg_exclusion(req: KGExclusionCreate, current_user: dict = Depends(auth.require_role("admin"))):
    """FR-30: Add entity exclusion and optionally delete existing nodes"""
    import sqlite3
    
    entity_name = req.entity_name.strip()
    if not entity_name:
        raise HTTPException(status_code=400, detail="Entity name is required")
        
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    try:
        # Add to exclusion list
        c.execute("INSERT INTO kg_exclusions (entity_name) VALUES (?)", (entity_name,))
        
        # Auto-cleanup: Delete any existing nodes with this exact label (case-insensitive)
        c.execute("SELECT id FROM kg_nodes WHERE LOWER(label) = LOWER(?)", (entity_name,))
        nodes_to_delete = [row[0] for row in c.fetchall()]
        
        deleted_nodes = len(nodes_to_delete)
        deleted_edges = 0
        
        if nodes_to_delete:
            placeholders = ",".join(["?"] * len(nodes_to_delete))
            # Delete connected edges
            c.execute(f"DELETE FROM kg_edges WHERE source_id IN ({placeholders}) OR target_id IN ({placeholders})", nodes_to_delete + nodes_to_delete)
            deleted_edges = c.rowcount
            # Delete the nodes
            c.execute(f"DELETE FROM kg_nodes WHERE id IN ({placeholders})", nodes_to_delete)
            
        conn.commit()
    except sqlite3.IntegrityError:
        conn.close()
        raise HTTPException(status_code=400, detail="Entity already in exclusion list")
    except Exception as e:
        conn.close()
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        conn.close()
        
    return {"message": "Success", "deleted_nodes": deleted_nodes, "deleted_edges": deleted_edges}

@app.delete("/api/admin/kg-exclusions/{exc_id}")
async def delete_kg_exclusion(exc_id: int, current_user: dict = Depends(auth.require_role("admin"))):
    """FR-30: Remove entity exclusion"""
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute("DELETE FROM kg_exclusions WHERE id = ?", (exc_id,))
    conn.commit()
    conn.close()
    return {"message": "Deleted successfully"}

# ── FR-31: System Monitoring (Insinyur TI) ──────────────────────────────────
import psutil
import shutil
import time
from datetime import datetime

@app.get("/api/engineer/health")
async def get_system_health(current_user: dict = Depends(auth.require_exact_role("insinyur ti"))):
    """FR-31: Fetch real-time system health metrics"""
    cpu_percent = psutil.cpu_percent(interval=0.5)
    mem = psutil.virtual_memory()
    disk = psutil.disk_usage('/')
    
    # DB File sizes
    db_metadata_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    db_users_path = os.path.join(BASE_DIR, "data", "users.db")
    
    metadata_size = os.path.getsize(db_metadata_path) if os.path.exists(db_metadata_path) else 0
    users_size = os.path.getsize(db_users_path) if os.path.exists(db_users_path) else 0

    return {
        "cpu": cpu_percent,
        "memory": {
            "total": mem.total,
            "used": mem.used,
            "percent": mem.percent
        },
        "disk": {
            "total": disk.total,
            "used": disk.used,
            "percent": disk.percent
        },
        "database": {
            "metadata_db_mb": round(metadata_size / (1024 * 1024), 2),
            "users_db_mb": round(users_size / (1024 * 1024), 2)
        },
        "uptime_seconds": int(time.time() - psutil.boot_time())
    }

@app.get("/api/engineer/queue")
async def get_processing_queue(current_user: dict = Depends(auth.require_exact_role("insinyur ti"))):
    """FR-31: Mock processing queue based on recent audit logs"""
    import sqlite3
    db_path = os.path.join(BASE_DIR, "data", "legal_metadata.db")
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute("SELECT user_id, action_type, resource_id, timestamp FROM audit_logs WHERE action_type IN ('UPLOAD_DOCUMENT', 'DELETE_DOCUMENT', 'REBUILD_GRAPH') ORDER BY timestamp DESC LIMIT 10")
    recent_tasks = [{"user": row[0], "action": row[1], "resource": row[2], "timestamp": row[3], "status": "COMPLETED"} for row in c.fetchall()]
    conn.close()
    
    return {"active_tasks": [], "recent_history": recent_tasks}

@app.get("/api/engineer/backups")
async def get_backups(current_user: dict = Depends(auth.require_exact_role("insinyur ti"))):
    """FR-31: List existing backups in the data directory"""
    data_dir = os.path.join(BASE_DIR, "data")
    backups = []
    if os.path.exists(data_dir):
        for f in os.listdir(data_dir):
            if f.endswith('.db') and 'backup' in f:
                path = os.path.join(data_dir, f)
                backups.append({
                    "filename": f,
                    "size_mb": round(os.path.getsize(path) / (1024 * 1024), 2),
                    "created_at": datetime.fromtimestamp(os.path.getctime(path)).isoformat()
                })
    return {"backups": sorted(backups, key=lambda x: x['created_at'], reverse=True)}

@app.post("/api/engineer/backup")
async def create_backup(current_user: dict = Depends(auth.require_exact_role("insinyur ti"))):
    """FR-31: Manually trigger SQLite database backups"""
    import sqlite3
    data_dir = os.path.join(BASE_DIR, "data")
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    dbs = ["legal_metadata.db", "users.db"]
    created_backups = []
    
    for db_name in dbs:
        src = os.path.join(data_dir, db_name)
        if os.path.exists(src):
            dst_name = db_name.replace('.db', f'_backup_{timestamp}.db')
            dst = os.path.join(data_dir, dst_name)
            
            # Use SQLite backup API for safe copy
            src_conn = sqlite3.connect(src)
            dst_conn = sqlite3.connect(dst)
            with dst_conn:
                src_conn.backup(dst_conn)
            dst_conn.close()
            src_conn.close()
            
            created_backups.append(dst_name)
            
    return {"message": "Backup successful", "files": created_backups}

if __name__ == "__main__":
    import uvicorn
    # Natively running on port 8000
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=False)
